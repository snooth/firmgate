"""AI Document Search — index text from portal files and answer via an OpenAI-compatible LLM."""

from __future__ import annotations

import fnmatch
import json
import logging
import os
import re
import ssl
from collections import defaultdict
import urllib.error
import urllib.request
from io import BytesIO
from pathlib import Path
from typing import Any

from flask import current_app, url_for
from sqlalchemy import inspect

from app import access
from app.extensions import db
from app.file_storage import absolute_path
from app.models import AiDocChunk, FileNode, FileVersion, User, utcnow
from app.settings import get_setting, set_setting

SETTING_AI_DOC_SEARCH = "ai_document_search"

log = logging.getLogger(__name__)

_MAX_FILE_BYTES = 25 * 1024 * 1024
_MAX_FILE_CHARS = 600_000
_CHUNK_CHARS = 2000
_MAX_CHUNKS_PER_FILE = 150
_INDEX_BATCH_FILES = 120
_SEARCH_TOP_K = 18
_EMBED_BATCH_SIZE = 48
_VECTOR_SCAN_LIMIT = 12000

_TEXT_SUFFIXES = frozenset(
    {
        ".txt",
        ".md",
        ".markdown",
        ".csv",
        ".tsv",
        ".log",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".css",
        ".scss",
        ".sql",
        ".sh",
        ".bash",
        ".zsh",
        ".yaml",
        ".yml",
        ".ini",
        ".cfg",
        ".conf",
        ".env",
        ".rb",
        ".go",
        ".rs",
        ".java",
        ".kt",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".php",
        ".vue",
        ".r",
        ".swift",
        ".toml",
        ".rst",
        ".tex",
        ".bat",
        ".ps1",
        ".pl",
        ".lua",
        ".vb",
        ".cs",
        ".m",
        ".mm",
        ".properties",
        ".gradle",
        ".dockerfile",
        ".makefile",
        ".gitignore",
        ".editorconfig",
    }
)

_INDEX_SUFFIXES = _TEXT_SUFFIXES | frozenset({".pdf", ".docx", ".xlsx", ".pptx", ".xlsm", ".pptm", ".docm"})


def _ensure_schema() -> None:
    try:
        insp = inspect(db.engine)
        if not insp.has_table("ai_doc_chunks"):
            db.create_all()
        else:
            cols = {c["name"] for c in insp.get_columns("ai_doc_chunks")}
            if "embedding_json" not in cols:
                from sqlalchemy import text

                with db.engine.begin() as conn:
                    conn.execute(text("ALTER TABLE ai_doc_chunks ADD COLUMN embedding_json TEXT"))
            if "index_kind" not in cols:
                from sqlalchemy import text

                with db.engine.begin() as conn:
                    conn.execute(
                        text(
                            "ALTER TABLE ai_doc_chunks ADD COLUMN index_kind VARCHAR(32) NOT NULL DEFAULT 'documents'"
                        )
                    )
    except Exception:
        try:
            db.create_all()
        except Exception:
            pass


def _settings_row() -> dict[str, Any]:
    raw = get_setting(SETTING_AI_DOC_SEARCH, default={}) or {}
    return raw if isinstance(raw, dict) else {}


def _setting_int(
    row_key: str,
    *,
    config_key: str | None = None,
    env_key: str | None = None,
    default: int,
    lo: int,
    hi: int,
) -> int:
    row = _settings_row()
    if row_key in row and row[row_key] is not None and str(row[row_key]).strip() != "":
        try:
            return max(lo, min(int(row[row_key]), hi))
        except (TypeError, ValueError):
            pass
    for src in (
        (current_app.config.get(config_key) if config_key else None),
        (os.environ.get(env_key) if env_key else None),
    ):
        if src is None or str(src).strip() == "":
            continue
        try:
            return max(lo, min(int(src), hi))
        except (TypeError, ValueError):
            continue
    return max(lo, min(default, hi))


def _parse_glob_list(raw: str | list | None) -> list[str]:
    if isinstance(raw, list):
        parts = [str(x).strip() for x in raw if str(x).strip()]
    else:
        parts = [p.strip() for p in re.split(r"[,;\n]+", (raw or "").strip()) if p.strip()]
    return parts[:40]


def index_scope_mode() -> str:
    mode = (_settings_row().get("index_scope_mode") or "all").strip().lower()
    return "folders" if mode == "folders" else "all"


def index_folder_ids() -> list[int]:
    raw = _settings_row().get("index_folder_ids") or []
    if not isinstance(raw, list):
        return []
    out: list[int] = []
    for x in raw:
        try:
            out.append(int(x))
        except (TypeError, ValueError):
            continue
    return out[:200]


def index_include_globs() -> list[str]:
    row = _settings_row()
    return _parse_glob_list(row.get("index_include_glob") or row.get("index_include_globs"))


def index_exclude_globs() -> list[str]:
    row = _settings_row()
    return _parse_glob_list(row.get("index_exclude_glob") or row.get("index_exclude_globs"))


def max_file_chars() -> int:
    return _setting_int(
        "max_file_chars",
        config_key="AI_DOC_MAX_FILE_CHARS",
        env_key="AI_DOC_MAX_FILE_CHARS",
        default=_MAX_FILE_CHARS,
        lo=10_000,
        hi=2_000_000,
    )


def embed_batch_size() -> int:
    return _setting_int(
        "embed_batch_size",
        config_key="AI_DOC_EMBED_BATCH",
        env_key="AI_DOC_EMBED_BATCH",
        default=_EMBED_BATCH_SIZE,
        lo=8,
        hi=128,
    )


def search_top_k_default() -> int:
    return _setting_int(
        "search_top_k",
        config_key="AI_DOC_SEARCH_TOP_K",
        env_key="AI_DOC_SEARCH_TOP_K",
        default=_SEARCH_TOP_K,
        lo=5,
        hi=50,
    )


def _collect_file_ids_under_folders(folder_ids: list[int]) -> set[int]:
    roots = {int(x) for x in folder_ids if x}
    if not roots:
        return set()
    rows = (
        db.session.query(FileNode.id, FileNode.parent_id, FileNode.is_folder)
        .filter(FileNode.deleted_at.is_(None))
        .all()
    )
    children: dict[int, list[int]] = defaultdict(list)
    file_ids: set[int] = set()
    folder_set = {fid for fid, _pid, is_folder in rows if is_folder}
    for fid, pid, is_folder in rows:
        if is_folder and pid is not None:
            children[int(pid)].append(int(fid))
    subtree = set(roots & folder_set)
    queue = list(subtree)
    while queue:
        fid = queue.pop()
        for child in children.get(fid, []):
            if child not in subtree:
                subtree.add(child)
                queue.append(child)
    for fid, pid, is_folder in rows:
        if not is_folder and pid is not None and int(pid) in subtree:
            file_ids.add(int(fid))
    return file_ids


def global_scoped_file_ids() -> set[int] | None:
    """``None`` = no folder restriction (all files may be indexed)."""
    if index_scope_mode() != "folders":
        return None
    return _collect_file_ids_under_folders(index_folder_ids())


def _node_passes_filename_filters(node: FileNode) -> bool:
    name = (node.name or "").strip()
    if not name:
        return False
    inc = index_include_globs()
    exc = index_exclude_globs()
    low = name.lower()
    if inc:
        if not any(fnmatch.fnmatch(low, p.lower()) for p in inc):
            return False
    if exc:
        if any(fnmatch.fnmatch(low, p.lower()) for p in exc):
            return False
    return True


def node_in_index_scope(node: FileNode) -> bool:
    if node.is_folder or node.deleted_at is not None:
        return False
    try:
        from app.enterprise.ai_policy_assistant import node_under_policy_folder

        if node_under_policy_folder(node):
            return False
    except Exception:
        pass
    if not is_indexable_filename(node.name):
        return False
    if not _node_passes_filename_filters(node):
        return False
    scoped = global_scoped_file_ids()
    if scoped is not None and node.id not in scoped:
        return False
    return True


def purge_chunks_outside_scope() -> int:
    _ensure_schema()
    scoped = global_scoped_file_ids()
    inc = index_include_globs()
    exc = index_exclude_globs()
    if scoped is None and not inc and not exc:
        return 0
    to_delete: list[int] = []
    for row in db.session.query(AiDocChunk.file_node_id).distinct().all():
        fid = int(row[0])
        node = db.session.get(FileNode, fid)
        if not node or node.deleted_at or node.is_folder:
            to_delete.append(fid)
            continue
        if not node_in_index_scope(node):
            to_delete.append(fid)
    if not to_delete:
        return 0
    deleted = (
        db.session.query(AiDocChunk)
        .filter(AiDocChunk.file_node_id.in_(to_delete))
        .delete(synchronize_session=False)
    )
    db.session.commit()
    return int(deleted or 0)


def list_index_folder_options(*, parent_id: int | None = None, q: str = "") -> list[dict[str, Any]]:
    qry = db.session.query(FileNode).filter(
        FileNode.deleted_at.is_(None),
        FileNode.is_folder.is_(True),
    )
    term = (q or "").strip()
    if term:
        qry = qry.filter(FileNode.name.ilike(f"%{term}%"))
        qry = qry.order_by(FileNode.name.asc()).limit(80)
    elif parent_id is not None:
        qry = qry.filter(FileNode.parent_id == int(parent_id))
        qry = qry.order_by(FileNode.name.asc()).limit(200)
    else:
        qry = qry.filter(FileNode.parent_id.is_(None))
        qry = qry.order_by(FileNode.name.asc()).limit(200)
    return [
        {
            "id": n.id,
            "name": n.name,
            "path": n.display_path(),
            "parent_id": n.parent_id,
        }
        for n in qry.all()
    ]


def resolve_index_folder_labels(folder_ids: list[int]) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for fid in folder_ids:
        node = db.session.get(FileNode, int(fid))
        if not node or node.deleted_at or not node.is_folder:
            continue
        out.append({"id": node.id, "path": node.display_path(), "name": node.name})
    return out


def indexing_settings_for_api() -> dict[str, Any]:
    folder_ids = index_folder_ids()
    return {
        "chunk_chars": _chunk_limits()[0],
        "max_chunks_per_file": _chunk_limits()[1],
        "max_file_chars": max_file_chars(),
        "embed_batch_size": embed_batch_size(),
        "search_top_k": search_top_k_default(),
        "index_batch_files": _index_batch_limit(),
        "index_scope_mode": index_scope_mode(),
        "index_folder_ids": folder_ids,
        "index_folders": resolve_index_folder_labels(folder_ids),
        "index_include_glob": ", ".join(index_include_globs()),
        "index_exclude_glob": ", ".join(index_exclude_globs()),
    }


def apply_indexing_settings_payload(payload: dict[str, Any], existing: dict[str, Any]) -> dict[str, Any]:
    """Merge indexing fields from admin PUT into settings dict."""
    merged = dict(existing)
    int_fields = {
        "chunk_chars": (400, 8000),
        "max_chunks_per_file": (10, 300),
        "max_file_chars": (10_000, 2_000_000),
        "embed_batch_size": (8, 128),
        "search_top_k": (5, 50),
        "index_batch_files": (10, 500),
    }
    for key, (lo, hi) in int_fields.items():
        if key not in payload:
            continue
        try:
            merged[key] = max(lo, min(int(payload[key]), hi))
        except (TypeError, ValueError):
            pass
    if "index_scope_mode" in payload:
        mode = (str(payload.get("index_scope_mode") or "all")).strip().lower()
        merged["index_scope_mode"] = "folders" if mode == "folders" else "all"
    if "index_folder_ids" in payload:
        raw = payload.get("index_folder_ids")
        ids: list[int] = []
        if isinstance(raw, list):
            for x in raw:
                try:
                    ids.append(int(x))
                except (TypeError, ValueError):
                    continue
        merged["index_folder_ids"] = ids[:200]
    for glob_key in ("index_include_glob", "index_exclude_glob"):
        if glob_key in payload:
            merged[glob_key] = str(payload.get(glob_key) or "").strip()[:2000]
    return merged


def llm_api_key() -> str:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_api_key as _key

    return _key(PRODUCT_DOCUMENT_SEARCH)


def llm_base_url() -> str:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_base_url as _url

    return _url(PRODUCT_DOCUMENT_SEARCH)


def llm_model() -> str:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_model as _model

    return _model(PRODUCT_DOCUMENT_SEARCH)


def llm_embedding_model() -> str:
    row = _settings_row()
    model = (row.get("embedding_model") or "").strip()
    if model:
        return model
    return (
        (current_app.config.get("AI_LLM_EMBEDDING_MODEL") or "")
        or (os.environ.get("AI_LLM_EMBEDDING_MODEL") or "text-embedding-3-small")
    ).strip()


def vector_index_version() -> int:
    row = _settings_row()
    try:
        return int(row.get("vector_index_version") or 0)
    except (TypeError, ValueError):
        return 0


def _bump_vector_index_version() -> int:
    row = _settings_row()
    try:
        n = int(row.get("vector_index_version") or 0) + 1
    except (TypeError, ValueError):
        n = 1
    row = dict(row)
    row["vector_index_version"] = n
    set_setting(SETTING_AI_DOC_SEARCH, row)
    return n


def _chunk_limits() -> tuple[int, int]:
    chars = _setting_int(
        "chunk_chars",
        config_key="AI_DOC_CHUNK_CHARS",
        env_key="AI_DOC_CHUNK_CHARS",
        default=_CHUNK_CHARS,
        lo=400,
        hi=8000,
    )
    max_chunks = _setting_int(
        "max_chunks_per_file",
        config_key="AI_DOC_MAX_CHUNKS",
        env_key="AI_DOC_MAX_CHUNKS",
        default=_MAX_CHUNKS_PER_FILE,
        lo=10,
        hi=300,
    )
    return chars, max_chunks


def llm_skip_tls_verify() -> bool:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_skip_tls_verify as _skip

    return _skip(PRODUCT_DOCUMENT_SEARCH)


def _llm_ssl_context() -> ssl.SSLContext:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, _llm_ssl_context as _ctx

    return _ctx(PRODUCT_DOCUMENT_SEARCH)


def llm_configured() -> bool:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_configured as _configured

    return _configured(PRODUCT_DOCUMENT_SEARCH)


def llm_settings_public() -> dict[str, Any]:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, llm_settings_public as _public

    return {
        **_public(PRODUCT_DOCUMENT_SEARCH),
        "embedding_model": llm_embedding_model(),
        "vector_index_version": vector_index_version(),
        **indexing_settings_for_api(),
    }


def _suffix(name: str) -> str:
    return Path(name or "").suffix.lower()


def is_indexable_filename(name: str) -> bool:
    suf = _suffix(name)
    if suf in _INDEX_SUFFIXES:
        return True
    if not suf and (name or "").lower() in ("makefile", "dockerfile", "license", "readme"):
        return True
    return False


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    if len(raw) > _MAX_FILE_BYTES:
        raw = raw[:_MAX_FILE_BYTES]
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for page in reader.pages[:200]:
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            parts.append(t)
    return "\n".join(parts)


def _docx_text(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    parts: list[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx_text(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets[:20]:
        parts.append(f"[Sheet: {sheet.title}]")
        row_count = 0
        for row in sheet.iter_rows(max_row=500, values_only=True):
            row_count += 1
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append("\t".join(cells))
        if row_count >= 500:
            parts.append("…")
    wb.close()
    return "\n".join(parts)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides[:80], start=1):
        slide_bits: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and (shape.text or "").strip():
                slide_bits.append(shape.text.strip())
        if slide_bits:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_bits))
    return "\n\n".join(parts)


def extract_text_from_bytes(name: str, data: bytes) -> str:
    if len(data) > _MAX_FILE_BYTES:
        data = data[:_MAX_FILE_BYTES]
    suf = _suffix(name)
    if suf == ".pdf":
        return _pdf_text(data)
    if suf == ".docx" or suf == ".docm":
        return _docx_text(data)
    if suf in (".xlsx", ".xlsm"):
        return _xlsx_text(data)
    if suf in (".pptx", ".pptm"):
        return _pptx_text(data)
    if suf in _TEXT_SUFFIXES or not suf:
        return _read_text_file_bytes(data)
    return ""


def _read_text_file_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text_from_path(name: str, path: Path) -> str:
    data = path.read_bytes()
    return extract_text_from_bytes(name, data)


def _chunk_text(text: str) -> list[str]:
    chunk_chars, max_chunks = _chunk_limits()
    text = re.sub(r"\r\n?", "\n", text or "")
    text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
    if not text:
        return []
    cap = max_file_chars()
    if len(text) > cap:
        text = text[:cap] + "\n… [truncated]"
    chunks: list[str] = []
    start = 0
    while start < len(text) and len(chunks) < max_chunks:
        end = min(len(text), start + chunk_chars)
        if end < len(text):
            break_at = text.rfind("\n\n", start, end)
            if break_at > start + chunk_chars // 3:
                end = break_at
            else:
                break_at = text.rfind("\n", start, end)
                if break_at > start + chunk_chars // 3:
                    end = break_at
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        start = end if end > start else start + chunk_chars
    return chunks


def _parse_embedding(raw: Any) -> list[float] | None:
    if raw is None:
        return None
    if isinstance(raw, str):
        s = raw.strip()
        if not s:
            return None
        try:
            data = json.loads(s)
        except json.JSONDecodeError:
            return None
    elif isinstance(raw, list):
        data = raw
    else:
        return None
    if not isinstance(data, list) or not data:
        return None
    try:
        return [float(x) for x in data]
    except (TypeError, ValueError):
        return None


def _serialize_embedding(vec: list[float]) -> str:
    return json.dumps(vec, separators=(",", ":"))


def _cosine_similarity(a: list[float], b: list[float]) -> float:
    if len(a) != len(b) or not a:
        return 0.0
    dot = 0.0
    na = 0.0
    nb = 0.0
    for x, y in zip(a, b):
        dot += x * y
        na += x * x
        nb += y * y
    if na <= 0.0 or nb <= 0.0:
        return 0.0
    return dot / ((na**0.5) * (nb**0.5))


def _llm_api_post(path_suffix: str, payload: dict[str, Any], *, timeout: int = 120) -> dict[str, Any]:
    from app.enterprise.ai_llm_http import effective_llm_api_key, openai_request_headers

    base = llm_base_url()
    api_key = llm_api_key()
    if not effective_llm_api_key(api_key, base):
        raise RuntimeError("AI API key is not configured.")
    url = f"{base}/{path_suffix.lstrip('/')}"
    body = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=body,
        headers=openai_request_headers(api_key, base),
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout, context=_llm_ssl_context()) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        err_body = e.read().decode("utf-8", errors="replace")[:500]
        raise RuntimeError(f"AI API request failed ({e.code}): {err_body}") from e
    except urllib.error.URLError as e:
        raise RuntimeError(f"Could not reach AI API: {e.reason}") from e


def _embed_texts(texts: list[str]) -> list[list[float]]:
    if not texts:
        return []
    model = llm_embedding_model()
    out: list[list[float]] = []
    batch_sz = embed_batch_size()
    for i in range(0, len(texts), batch_sz):
        batch = texts[i : i + batch_sz]
        data = _llm_api_post("embeddings", {"model": model, "input": batch}, timeout=180)
        items = data.get("data")
        if not isinstance(items, list):
            raise RuntimeError("Unexpected embeddings API response.")
        ordered = sorted(items, key=lambda x: int(x.get("index", 0)))
        for item in ordered:
            emb = item.get("embedding")
            vec = _parse_embedding(emb)
            if not vec:
                raise RuntimeError("Invalid embedding vector in API response.")
            out.append(vec)
    if len(out) != len(texts):
        raise RuntimeError("Embeddings API returned an unexpected number of vectors.")
    return out


def _file_missing_embeddings(file_node_id: int, index_kind: str = "documents") -> bool:
    kind = (index_kind or "documents").strip().lower() or "documents"
    rows = db.session.query(AiDocChunk.embedding_json).filter_by(file_node_id=file_node_id, index_kind=kind).all()
    if not rows:
        return True
    return any(not (r[0] or "").strip() for r in rows)


def _latest_version(node: FileNode) -> FileVersion | None:
    return (
        db.session.query(FileVersion)
        .filter_by(file_node_id=node.id)
        .order_by(FileVersion.version_number.desc())
        .first()
    )


def _user_can_read(user: User, node: FileNode) -> bool:
    try:
        ok, _ = access.can_access_node(user, node, "read")
        return bool(ok)
    except Exception:
        return False


def index_stats() -> dict[str, int]:
    _ensure_schema()
    try:
        files = db.session.query(AiDocChunk.file_node_id).distinct().count()
        chunks = db.session.query(AiDocChunk).count()
    except Exception:
        files = 0
        chunks = 0
    return {"indexed_files": files, "indexed_chunks": chunks}


def index_stats_for_user(user: User) -> dict[str, int]:
    """Counts indexed chunks for files this user can read (for UI)."""
    _ensure_schema()
    file_ids: set[int] = set()
    chunk_count = 0
    try:
        rows = db.session.query(AiDocChunk).all()
    except Exception:
        return {"indexed_files": 0, "indexed_chunks": 0}
    for row in rows:
        node = db.session.get(FileNode, row.file_node_id)
        if not node or node.deleted_at or not _user_can_read(user, node):
            continue
        file_ids.add(row.file_node_id)
        chunk_count += 1
    return {"indexed_files": len(file_ids), "indexed_chunks": chunk_count}


def _delete_chunks(file_node_id: int) -> None:
    db.session.query(AiDocChunk).filter_by(file_node_id=file_node_id).delete()


def index_file_node(user: User, node: FileNode, *, force: bool = False, index_kind: str = "documents") -> bool:
    """Index one file if readable. Returns True if indexed or already current."""
    _ensure_schema()
    if node.is_folder or node.deleted_at is not None:
        return False
    kind = (index_kind or "documents").strip().lower() or "documents"
    if kind == "documents" and not node_in_index_scope(node):
        return False
    if not _user_can_read(user, node):
        return False
    ver = _latest_version(node)
    if not ver:
        return False
    existing = (
        db.session.query(AiDocChunk)
        .filter_by(file_node_id=node.id, index_kind=kind)
        .limit(1)
        .first()
    )
    if (
        existing
        and existing.file_sha256 == ver.sha256
        and not force
        and not _file_missing_embeddings(node.id, kind)
    ):
        return True
    try:
        path = absolute_path(ver.storage_relpath)
        text = extract_text_from_path(node.name, path)
    except Exception as exc:
        log.debug("ai index skip %s: %s", node.name, exc)
        return False
    chunks = _chunk_text(text)
    if not chunks:
        db.session.query(AiDocChunk).filter_by(file_node_id=node.id, index_kind=kind).delete()
        db.session.commit()
        return False
    path_label = node.display_path()
    embeddings: list[list[float]] = []
    if llm_configured():
        try:
            embeddings = _embed_texts(chunks)
        except Exception as exc:
            log.warning("AI embeddings failed for %s: %s", node.name, exc)
    db.session.query(AiDocChunk).filter_by(file_node_id=node.id, index_kind=kind).delete()
    for i, body in enumerate(chunks):
        emb_json = None
        if i < len(embeddings):
            emb_json = _serialize_embedding(embeddings[i])
        db.session.add(
            AiDocChunk(
                file_node_id=node.id,
                file_sha256=ver.sha256,
                chunk_index=i,
                path_label=path_label[:1024],
                body_text=body,
                embedding_json=emb_json,
                index_kind=kind,
                indexed_at=utcnow(),
            )
        )
    db.session.commit()
    if embeddings:
        _bump_vector_index_version()
    return True


def _index_batch_limit() -> int:
    return _setting_int(
        "index_batch_files",
        config_key="AI_DOC_INDEX_BATCH",
        env_key="AI_DOC_INDEX_BATCH",
        default=_INDEX_BATCH_FILES,
        lo=10,
        hi=500,
    )


def sync_index_for_user(user: User, *, limit: int | None = None) -> dict[str, int]:
    """Index up to ``limit`` stale/missing files the user can read."""
    _ensure_schema()
    batch = limit if limit is not None else _index_batch_limit()
    nodes = (
        db.session.query(FileNode)
        .filter(FileNode.deleted_at.is_(None), FileNode.is_folder.is_(False))
        .order_by(FileNode.updated_at.desc())
        .limit(5000)
        .all()
    )
    indexed = 0
    scanned = 0
    skipped_current = 0
    skipped_no_access = 0
    skipped_not_supported = 0
    skipped_no_version = 0
    skipped_empty = 0
    skipped_error = 0
    skipped_out_of_scope = 0
    for node in nodes:
        if indexed >= batch:
            break
        if scanned >= batch * 8 and indexed >= batch:
            break
        if not is_indexable_filename(node.name):
            skipped_not_supported += 1
            continue
        if not node_in_index_scope(node):
            skipped_out_of_scope += 1
            continue
        if not _user_can_read(user, node):
            skipped_no_access += 1
            continue
        scanned += 1
        ver = _latest_version(node)
        if not ver:
            skipped_no_version += 1
            continue
        row = (
            db.session.query(AiDocChunk.file_sha256)
            .filter_by(file_node_id=node.id)
            .limit(1)
            .first()
        )
        if row and row[0] == ver.sha256 and not _file_missing_embeddings(node.id):
            skipped_current += 1
            continue
        had_chunks = (
            db.session.query(AiDocChunk.id).filter_by(file_node_id=node.id).limit(1).first()
            is not None
        )
        if index_file_node(user, node):
            indexed += 1
        elif not had_chunks:
            skipped_empty += 1
        else:
            skipped_error += 1
    return {
        "indexed_now": indexed,
        "scanned": scanned,
        "skipped_current": skipped_current,
        "skipped_no_access": skipped_no_access,
        "skipped_not_supported": skipped_not_supported,
        "skipped_no_version": skipped_no_version,
        "skipped_empty": skipped_empty,
        "skipped_error": skipped_error,
        "skipped_out_of_scope": skipped_out_of_scope,
        "batch_limit": batch,
        "index_scope_mode": index_scope_mode(),
    }


def _score_chunk(body: str, terms: list[str]) -> int:
    low = body.lower()
    score = 0
    for t in terms:
        if t in low:
            score += low.count(t) * 2 + 3
    return score


def search_chunks(
    user: User,
    query: str,
    *,
    limit: int | None = None,
    index_kind: str | None = "documents",
) -> list[dict[str, Any]]:
    if limit is None:
        limit = search_top_k_default()
    _ensure_schema()
    q = (query or "").strip()
    if not q:
        return []
    kind = (index_kind or "").strip().lower() if index_kind else ""
    terms = [t for t in re.findall(r"[a-zA-Z0-9_]{2,}", q.lower())[:12]]
    if not terms:
        terms = [q.lower()[:64]]

    query_vec: list[float] | None = None
    if llm_configured():
        try:
            query_vec = _embed_texts([q[:8000]])[0]
        except Exception as exc:
            log.debug("query embedding failed: %s", exc)

    rows = (
        db.session.query(AiDocChunk)
        .order_by(AiDocChunk.indexed_at.desc())
        .limit(_VECTOR_SCAN_LIMIT)
        .all()
    )
    scored: list[tuple[float, AiDocChunk]] = []
    for row in rows:
        row_kind = (getattr(row, "index_kind", None) or "documents").strip().lower()
        if kind and row_kind != kind:
            continue
        node = db.session.get(FileNode, row.file_node_id)
        if not node or node.deleted_at or not _user_can_read(user, node):
            continue
        kw = float(_score_chunk(row.body_text, terms))
        vec = _parse_embedding(row.embedding_json)
        if query_vec and vec:
            sim = _cosine_similarity(query_vec, vec)
            score = sim * 100.0 + kw
            if score < 0.12 and kw <= 0:
                continue
        elif kw <= 0:
            continue
        else:
            score = kw
        scored.append((score, row))
    scored.sort(key=lambda x: (-x[0], x[1].file_node_id, x[1].chunk_index))
    hits: list[dict[str, Any]] = []
    for score, row in scored:
        if len(hits) >= limit:
            break
        node = db.session.get(FileNode, row.file_node_id)
        if not node:
            continue
        snippet = row.body_text.replace("\n", " ")
        if len(snippet) > 320:
            snippet = snippet[:317] + "…"
        hits.append(
            {
                "file_node_id": row.file_node_id,
                "path": row.path_label,
                "snippet": snippet,
                "score": int(score * 100),
                "url": url_for(
                    "intranet.documents_page",
                    parent_id=node.parent_id,
                    select_id=row.file_node_id,
                ),
            }
        )
        if len(hits) >= limit:
            break
    return hits


def _build_context(hits: list[dict[str, Any]]) -> str:
    if not hits:
        return "(No matching document excerpts were found in the index.)"
    parts: list[str] = []
    for i, h in enumerate(hits, start=1):
        parts.append(
            f"[Source {i}] {h.get('path') or 'Document'}\n"
            f"URL path: documents (file id {h.get('file_node_id')})\n"
            f"{h.get('snippet') or ''}"
        )
    return "\n\n".join(parts)


def _chat_completion(messages: list[dict[str, str]]) -> str:
    from app.enterprise.ai_llm_settings import PRODUCT_DOCUMENT_SEARCH, chat_completion

    return chat_completion(PRODUCT_DOCUMENT_SEARCH, messages)


def _document_search_system_prompt() -> str:
    return (
        "You are a document assistant for an organisation's private intranet. "
        "Answer using ONLY the document excerpts provided. "
        "If the excerpts do not contain enough information, say so clearly and suggest what to search for. "
        "Cite sources as [Source N] matching the excerpt labels. "
        "Do not invent file names or policy details. Be concise and practical."
    )


def prepare_answer_messages(
    user: User,
    question: str,
    *,
    history: list[dict[str, str]] | None = None,
) -> dict[str, Any]:
    """Build LLM messages and retrieval hits without calling the model."""
    question = (question or "").strip()
    if not question:
        raise ValueError("Question is required.")
    if len(question) > 8000:
        raise ValueError("Question is too long.")
    sync_index_for_user(user)
    hits = search_chunks(user, question)
    context = _build_context(hits)
    user_content = f"Document excerpts:\n\n{context}\n\nUser question: {question}"
    messages: list[dict[str, str]] = [{"role": "system", "content": _document_search_system_prompt()}]
    if history:
        for m in history[-8:]:
            role = (m.get("role") or "").strip().lower()
            content = (m.get("content") or "").strip()
            if role in ("user", "assistant") and content:
                messages.append({"role": role, "content": content[:6000]})
    messages.append({"role": "user", "content": user_content})
    return {
        "messages": messages,
        "sources": hits,
        "stats": index_stats_for_user(user),
    }


def iter_document_search_deltas(
    messages: list[dict[str, str]],
    *,
    temperature: float | None = None,
    max_tokens: int | None = None,
    timeout: int = 300,
):
    from app.enterprise.ai_llm_settings import (
        PRODUCT_DOCUMENT_SEARCH,
        iter_chat_completion_deltas,
        llm_max_tokens,
        llm_temperature,
    )

    return iter_chat_completion_deltas(
        PRODUCT_DOCUMENT_SEARCH,
        messages,
        temperature=temperature if temperature is not None else llm_temperature(PRODUCT_DOCUMENT_SEARCH),
        max_tokens=max_tokens if max_tokens is not None else llm_max_tokens(PRODUCT_DOCUMENT_SEARCH),
        timeout=timeout,
    )


def answer_question(
    user: User,
    question: str,
    *,
    history: list[dict[str, str]] | None = None,
) -> dict[str, Any]:
    prep = prepare_answer_messages(user, question, history=history)
    answer = _chat_completion(prep["messages"])
    return {
        "answer": answer,
        "sources": prep["sources"],
        "stats": prep["stats"],
    }

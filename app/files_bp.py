from __future__ import annotations

import os
import zipfile
import mimetypes
from io import BytesIO
from typing import Any

from flask import Blueprint, abort, jsonify, redirect, render_template, request, send_file, url_for
from flask_login import current_user, login_required
from sqlalchemy.exc import SQLAlchemyError
from sqlalchemy import and_, func, or_, select

from app import access
from app import files_workspace
from app import rbac
from app.audit_service import validate_deletion_justification, write_audit
from app.extensions import db
from app.file_storage import absolute_path, copy_blob_to_new_version, store_stream_and_digest
from app.models import (
    AuditLog,
    FileComment,
    FileNode,
    FileShare,
    FileVersion,
    Group,
    NodeFavorite,
    NodeGroupShare,
    NodeRoleShare,
    NodeUserShare,
    Role,
    User,
    user_roles,
    utcnow,
)
from app.file_lock_service import (
    acquire_lock,
    assert_can_edit_locked_node,
    attach_lock_fields,
    get_lock_for_node,
    lock_payload,
    release_lock,
)
from app.file_edit_session_service import (
    attach_edit_fields,
    edit_sessions_map,
    release_edit_session,
    touch_edit_session,
)
from datetime import timedelta
from email import policy
from email.parser import BytesParser

bp = Blueprint("files", __name__, url_prefix="/files")

def _next_available_name(parent_id: int | None, desired: str) -> str:
    """Return desired or desired (n) (preserving file extension) if taken in parent."""
    desired = (desired or "").strip() or "item"
    if "/" in desired or "\\" in desired:
        desired = desired.replace("/", "_").replace("\\", "_")

    base, ext = os.path.splitext(desired)
    base = base.strip() or "item"

    def exists(name: str) -> bool:
        q = FileNode.query.filter_by(parent_id=parent_id, name=name).filter(FileNode.deleted_at.is_(None))
        return q.first() is not None

    if not exists(desired):
        return desired

    i = 1
    while i < 10_000:
        cand = f"{base} ({i}){ext}"
        if not exists(cand):
            return cand
        i += 1
    # extremely unlikely
    return f"{base} ({utcnow().timestamp():.0f}){ext}"


def _audit(action: str, resource_type: str | None, resource_id: str | None, success: bool, details: dict | None = None):
    write_audit(
        user_id=current_user.id,
        username=current_user.username,
        action=action,
        resource_type=resource_type,
        resource_id=resource_id,
        success=success,
        details=details,
    )


def _user_may_create_folder(user: User, parent: FileNode | None) -> bool:
    """Creating a subfolder under ``parent`` (non-root). Requires ``files.create_folders`` (owners included)."""
    if rbac.user_has_permission(user, rbac.PERMISSION_ADMIN):
        return True
    if rbac.user_has_permission(user, rbac.PERMISSION_FILES_ADMIN):
        return True
    if parent is not None and files_workspace.is_users_container_folder(parent):
        return False
    if parent is not None and parent.owner_id == user.id:
        return rbac.user_has_permission(user, rbac.PERMISSION_FILES_CREATE_FOLDERS)
    return rbac.user_has_permission(user, rbac.PERMISSION_FILES_CREATE_FOLDERS)


def _node_or_404(node_id: int) -> FileNode:
    node = db.session.get(FileNode, node_id)
    if not node:
        abort(404)
    return node


def _shows_users_root_display_name(user: User) -> bool:
    """Portal or files admins see the friendlier ``Users`` root label in the file browser."""
    return rbac.user_has_permission(user, rbac.PERMISSION_FILES_ADMIN) or rbac.user_has_permission(
        user, rbac.PERMISSION_ADMIN
    )


def _apply_users_root_display_breadcrumb(crumbs: list[dict[str, Any]]) -> None:
    if not _shows_users_root_display_name(current_user):
        return
    ur = files_workspace.find_users_root_folder()
    if not ur:
        return
    for c in crumbs:
        if c.get("id") == ur.id:
            c["name"] = files_workspace.USERS_ROOT_DISPLAY_NAME
            break


def _breadcrumb_chain(node: FileNode) -> list[dict[str, Any]]:
    parts: list[dict[str, Any]] = []
    cur: FileNode | None = node
    while cur is not None:
        parts.append({"id": cur.id, "name": cur.name})
        cur = cur.parent
    parts.reverse()
    return parts


def _profile_documents_anchor_ids(user: User) -> set[int]:
    """Folder ids hidden from Documents breadcrumbs (Users container, workspace, profile home)."""
    home = _default_home_folder_for_user(user.id)
    if not home:
        return set()
    ids: set[int] = set()
    cur: FileNode | None = home
    while cur is not None:
        ids.add(cur.id)
        if files_workspace.is_users_container_folder(cur):
            break
        cur = cur.parent
    return ids


def _documents_breadcrumb(user: User, node: FileNode | None) -> list[dict[str, Any]]:
    """All files = profile home; omit system path prefixes (Users Folder / username / Home)."""
    if node is None:
        return [{"id": None, "name": "All files"}]
    anchors = _profile_documents_anchor_ids(user)
    chain = [c for c in _breadcrumb_chain(node) if c.get("id") not in anchors]
    return [{"id": None, "name": "All files"}] + chain


def _admin_documents_breadcrumb(node: FileNode | None) -> list[dict[str, Any]]:
    """Full tree trail for files-admin browsing (Users Folder, all workspaces, etc.)."""
    if node is None:
        return [{"id": None, "name": "All users"}]
    crumbs = [{"id": None, "name": "All users"}] + _breadcrumb_chain(node)
    _apply_users_root_display_breadcrumb(crumbs)
    return crumbs



def _serialize_node_for_lister(
    n: FileNode,
    *,
    include_children_count: bool = False,
    shared_out_ids: set[int] | None = None,
    include_owner_username: bool = False,
) -> dict[str, Any]:
    batch = _serialize_nodes_batch(
        [n],
        shared_out_ids=shared_out_ids,
        include_owner_username=include_owner_username,
    )
    return batch[0] if batch else _serialize_node(n, include_children_count=include_children_count)


def _current_versions_map(file_node_ids: list[int]) -> dict[int, FileVersion]:
    if not file_node_ids:
        return {}
    rows = (
        db.session.query(FileVersion)
        .filter(FileVersion.file_node_id.in_(file_node_ids), FileVersion.is_current.is_(True))
        .all()
    )
    return {int(fv.file_node_id): fv for fv in rows}


def _shared_out_owned_node_ids(owner_id: int) -> set[int]:
    """Node IDs owned by this user that have any internal or public share."""
    owned = select(FileNode.id).where(
        FileNode.owner_id == owner_id,
        FileNode.deleted_at.is_(None),
    )
    out: set[int] = set()
    for model in (NodeUserShare, NodeGroupShare, NodeRoleShare, FileShare):
        rows = db.session.scalars(
            select(model.file_node_id).where(model.file_node_id.in_(owned)).distinct()
        ).all()
        out |= {int(i) for i in rows}
    return out


def _serialize_nodes_batch(
    nodes: list[FileNode],
    *,
    shared_out_ids: set[int] | None = None,
    include_owner_username: bool = False,
    attach_locks: bool = True,
) -> list[dict[str, Any]]:
    """Serialize many nodes with batched version, owner, and lock queries (list views)."""
    if not nodes:
        return []
    file_ids = [n.id for n in nodes if not n.is_folder]
    versions = _current_versions_map(file_ids)
    owner_names: dict[int, str] = {}
    if include_owner_username:
        owner_ids = {n.owner_id for n in nodes}
        if owner_ids:
            for u in db.session.query(User).filter(User.id.in_(owner_ids)).all():
                owner_names[u.id] = (u.username or u.email or "")

    items: list[dict[str, Any]] = []
    for n in nodes:
        data: dict[str, Any] = {
            "id": n.id,
            "name": n.name,
            "is_folder": n.is_folder,
            "parent_id": n.parent_id,
            "owner_id": n.owner_id,
            "created_at": n.created_at.isoformat() if n.created_at else None,
            "updated_at": n.updated_at.isoformat() if n.updated_at else None,
            "path_key": n.path_key,
            "attributes": n.attributes or {},
        }
        if not n.is_folder:
            fv = versions.get(n.id)
            if fv:
                data["size_bytes"] = fv.size_bytes
                data["mime_type"] = fv.mime_type
                data["version_number"] = fv.version_number
                data["sha256"] = fv.sha256
        if shared_out_ids is not None:
            data["shared_out"] = n.id in shared_out_ids
        if include_owner_username:
            data["owner_username"] = owner_names.get(n.owner_id, "")
        if files_workspace.is_users_container_folder(n) and _shows_users_root_display_name(current_user):
            data["name"] = files_workspace.USERS_ROOT_DISPLAY_NAME
        items.append(data)
    if attach_locks:
        attach_lock_fields(items)
    attach_edit_fields(items)
    return items


def _path_key_for(node: FileNode) -> str:
    if node.parent_id is None:
        return "/" + node.name
    parent = node.parent
    if not parent:
        return "/" + node.name
    base = (parent.path_key or _path_key_for(parent)).rstrip("/")
    return base + "/" + node.name


def _collect_subtree(root: FileNode) -> list[FileNode]:
    out: list[FileNode] = []
    stack = [root]
    while stack:
        n = stack.pop()
        out.append(n)
        if n.is_folder:
            for ch in n.children:
                stack.append(ch)
    return out


def _default_home_folder_for_user(user_id: int) -> FileNode | None:
    """Default upload/list folder: ``Users/<username>/`` (and legacy ``Home`` under it when present)."""
    return files_workspace.default_document_home_for_user(user_id)


def _is_files_tree_admin(user: User) -> bool:
    return rbac.user_has_permission(user, rbac.PERMISSION_FILES_ADMIN) or rbac.user_has_permission(
        user, rbac.PERMISSION_ADMIN
    )


def _is_portal_admin(user: User) -> bool:
    """Portal administrator (admin.all) — may release any file lock."""
    return rbac.user_has_permission(user, rbac.PERMISSION_ADMIN)


def _resolve_parent_for_create(parent_id: int | None, user: User) -> FileNode | None:
    """
    Resolve target folder for mkdir / office create when the UI sends a virtual root (parent_id null).
    Non-admins always create under their profile home folder.
    """
    if parent_id is not None:
        return _node_or_404(int(parent_id))
    if _is_files_tree_admin(user):
        return None
    parent = _default_home_folder_for_user(user.id)
    if not parent:
        return None
    return parent


def _delete_versions_for_node(node: FileNode) -> None:
    if node.is_folder:
        return
    for v in node.versions.all():
        relpath = v.storage_relpath
        db.session.delete(v)
        db.session.flush()
        remaining = db.session.query(FileVersion).filter_by(storage_relpath=relpath).count()
        if remaining == 0:
            try:
                p = absolute_path(relpath)
                if p.exists():
                    p.unlink()
            except (OSError, ValueError):
                pass


def _delete_shares_for_node_id(node_id: int) -> None:
    """Remove share rows so FKs do not block file_nodes delete (SQLAlchemy would null FK otherwise)."""
    NodeUserShare.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)
    NodeGroupShare.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)
    NodeRoleShare.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)
    FileShare.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)


def _delete_related_for_node_id(node_id: int) -> None:
    """Remove rows that FK to file_nodes so folder deletes don't fail."""
    FileComment.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)
    NodeFavorite.query.filter_by(file_node_id=node_id).delete(synchronize_session=False)


def _delete_node_recursive(node: FileNode) -> None:
    if node.is_folder:
        for ch in node.children.all():
            _delete_node_recursive(ch)
    else:
        _delete_versions_for_node(node)
    _delete_shares_for_node_id(node.id)
    _delete_related_for_node_id(node.id)
    db.session.delete(node)


def _soft_delete_node_recursive(node: FileNode, *, deleted_by_id: int) -> int:
    """Soft delete node + descendants into Recycle Bin. Returns count marked."""
    ts = utcnow()
    count = 0
    for n in _collect_subtree(node):
        if n.deleted_at is None:
            n.original_parent_id = n.parent_id
        n.deleted_at = ts
        n.deleted_by_id = deleted_by_id
        # Revoke any shares/public links so deleted content is no longer accessible.
        _delete_shares_for_node_id(n.id)
        if not n.is_folder:
            lk = get_lock_for_node(n.id)
            if lk:
                db.session.delete(lk)
        db.session.add(n)
        count += 1
    db.session.flush()
    return count


@bp.route("/")
@login_required
def browser():
    from app.community_edition import abort_if_module_locked

    abort_if_module_locked("documents")
    if not access.can_list_files(current_user):
        _audit("files.ui.denied", "page", "browser", False, {})
        abort(403)
    from app.document_editor_settings import files_template_context

    return render_template(
        "files.html",
        files_tree_admin=_is_files_tree_admin(current_user),
        portal_admin=_is_portal_admin(current_user),
        **files_template_context(),
    )


def _documents_root_list_payload(user: User) -> dict[str, Any]:
    """Virtual All files root: contents of the user's profile home folder."""
    home_container = _default_home_folder_for_user(user.id)
    items: list[dict[str, Any]] = []
    default_parent = home_container
    if home_container and home_container.is_folder:
        children = (
            home_container.children.filter(
                FileNode.deleted_at.is_(None),
                FileNode.owner_id == user.id,
            )
            .order_by(FileNode.is_folder.desc(), FileNode.name)
            .all()
        )
        visible = [c for c in children if access.documents_listing_includes_node(user, c)]
        child_ids = [c.id for c in visible]
        shared_out_ids = _shared_out_node_ids(child_ids) if child_ids else set()
        items = _serialize_nodes_batch(visible, shared_out_ids=shared_out_ids)

    shared_with_me: list[dict[str, Any]] = []
    seen_share_nodes: set[int] = set()
    _append_shared_with_me_entries(user, shared_with_me, seen_share_nodes)
    shared_with_me.sort(key=lambda x: (not x["is_folder"], x["name"].lower()))

    shared_by_me: list[dict[str, Any]] = []
    out_ids = _shared_out_owned_node_ids(user.id)
    if out_ids:
        shared_nodes = (
            FileNode.query.filter(FileNode.id.in_(out_ids), FileNode.deleted_at.is_(None))
            .order_by(FileNode.is_folder.desc(), FileNode.name)
            .all()
        )
        visible_shared = [sn for sn in shared_nodes if access.documents_listing_includes_node(user, sn)]
        shared_by_me = _serialize_nodes_batch(visible_shared)
        for d in shared_by_me:
            d["shared_out"] = True
    shared_by_me.sort(key=lambda x: (not x["is_folder"], (x.get("name") or "").lower()))

    _audit("files.list", "folder", "root", True, {"count": len(items)})
    return {
        "parent": None,
        "items": items,
        "breadcrumb": _documents_breadcrumb(user, None),
        "shared_with_me": shared_with_me,
        "shared_by_me": shared_by_me,
        "default_parent_id": (default_parent.id if default_parent else None),
        "profile_home_id": (default_parent.id if default_parent else None),
        "list_scope": "profile",
    }


def _documents_admin_root_list_payload(user: User) -> dict[str, Any]:
    """Files-admin view: every top-level folder in the file tree (Users, Security Training, etc.)."""
    roots = (
        FileNode.query.filter(FileNode.parent_id.is_(None), FileNode.deleted_at.is_(None))
        .order_by(FileNode.is_folder.desc(), FileNode.name)
        .all()
    )
    root_ids = [n.id for n in roots]
    shared_out_ids = _shared_out_node_ids(root_ids)
    visible = [
        n
        for n in roots
        if access.documents_listing_includes_node(user, n)
        and access.can_access_node(user, n, "read")[0]
    ]
    items = _serialize_nodes_batch(
        visible,
        shared_out_ids=shared_out_ids,
        include_owner_username=True,
    )
    _audit("files.list.admin", "folder", "root", True, {"count": len(items)})
    return {
        "parent": None,
        "items": items,
        "breadcrumb": _admin_documents_breadcrumb(None),
        "shared_with_me": [],
        "shared_by_me": [],
        "default_parent_id": None,
        "profile_home_id": None,
        "list_scope": "admin",
    }


@bp.route("/api/list", methods=["GET"])
@login_required
def api_list():
    if not access.can_list_files(current_user):
        _audit("files.list.denied", None, None, False, {})
        return jsonify({"error": "forbidden"}), 403
    scope = (request.args.get("scope") or "profile").strip().lower()
    admin_scope = scope == "admin"
    if admin_scope and not _is_files_tree_admin(current_user):
        return jsonify({"error": "forbidden", "reason": "files admin only"}), 403

    parent_id = request.args.get("parent_id", type=int)
    if parent_id is None:
        if admin_scope:
            return jsonify(_documents_admin_root_list_payload(current_user))
        return jsonify(_documents_root_list_payload(current_user))

    parent = _node_or_404(parent_id)
    profile_home = _default_home_folder_for_user(current_user.id)
    if not admin_scope and (
        (profile_home and parent.id == profile_home.id) or files_workspace.is_users_container_folder(parent)
    ):
        return jsonify(_documents_root_list_payload(current_user))
    if not parent.is_folder:
        return jsonify({"error": "not a folder"}), 400
    if parent.deleted_at is not None:
        return jsonify({"error": "not found"}), 404
    ok, reason = access.can_access_node(current_user, parent, "read")
    if not ok:
        _audit("files.list.denied", "folder", str(parent.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if not access.documents_listing_includes_node(current_user, parent):
        _audit("files.list.denied", "folder", str(parent.id), False, {"reason": "documents listing"})
        return jsonify({"error": "forbidden", "reason": "not available in Documents"}), 403

    children = (
        parent.children.filter(FileNode.deleted_at.is_(None))
        .order_by(FileNode.is_folder.desc(), FileNode.name)
        .all()
    )
    if files_workspace.is_users_container_folder(parent) and not admin_scope:
        children = files_workspace.filter_users_root_children_for_lister(current_user, list(children))
    child_ids = [c.id for c in children]
    shared_out_ids = _shared_out_node_ids(child_ids) if child_ids else set()
    visible: list[FileNode] = []
    for c in children:
        okc, _ = access.can_access_node(current_user, c, "read")
        if not okc:
            continue
        if not access.documents_listing_includes_node(current_user, c):
            continue
        visible.append(c)
    items = _serialize_nodes_batch(
        visible,
        shared_out_ids=shared_out_ids,
        include_owner_username=admin_scope,
    )
    # Treat navigating into a folder as "opened" (single-click should not spam activity).
    _audit("files.open", "folder", str(parent.id), True, {"count": len(items), "scope": scope})
    crumbs = _admin_documents_breadcrumb(parent) if admin_scope else _documents_breadcrumb(current_user, parent)
    return jsonify(
        {
            "parent": _serialize_node(parent),
            "items": items,
            "breadcrumb": crumbs,
            "shared_with_me": [],
            "profile_home_id": (profile_home.id if profile_home and not admin_scope else None),
            "list_scope": ("admin" if admin_scope else "profile"),
        }
    )


def _escape_ilike_pattern(q: str) -> str:
    """Escape SQL LIKE wildcards for safe ILIKE matching."""
    s = (q or "").replace("\\", "\\\\").replace("%", "\\%").replace("_", "\\_")
    return s


@bp.route("/api/search", methods=["GET"])
@login_required
def api_search():
    """Search files and folders by item name and full path (path_key). Space-separated words must all match."""
    if not access.can_list_files(current_user):
        _audit("files.search.denied", None, None, False, {})
        return jsonify({"error": "forbidden"}), 403

    raw = (request.args.get("q") or "").strip()
    tokens = [t for t in raw.split() if t]
    if not tokens:
        return jsonify({"error": "q required"}), 400
    if len(raw) > 500:
        return jsonify({"error": "query too long"}), 400
    for tok in tokens:
        if len(tok) > 200:
            return jsonify({"error": "token too long"}), 400

    token_conds: list[Any] = []
    for tok in tokens:
        pat = f"%{_escape_ilike_pattern(tok)}%"
        token_conds.append(
            or_(
                FileNode.name.ilike(pat, escape="\\"),
                func.coalesce(FileNode.path_key, "").ilike(pat, escape="\\"),
            )
        )
    name_or_path_match = and_(*token_conds)

    is_admin = rbac.user_has_permission(current_user, rbac.PERMISSION_FILES_ADMIN) or rbac.user_has_permission(
        current_user, rbac.PERMISSION_ADMIN
    )

    shared_ids_sq = (
        select(NodeUserShare.file_node_id)
        .where(NodeUserShare.shared_with_user_id == current_user.id)
        .distinct()
    )
    share_scope: list[Any] = [FileNode.owner_id == current_user.id, FileNode.id.in_(shared_ids_sq)]
    for prefix in _group_role_share_path_prefixes(current_user):
        esc = _escape_ilike_pattern(prefix)
        share_scope.append(FileNode.path_key == prefix)
        share_scope.append(FileNode.path_key.like(f"{esc}/%", escape="\\"))

    base_filters = [FileNode.deleted_at.is_(None), name_or_path_match]
    if is_admin:
        scope = FileNode.query.filter(*base_filters)
    else:
        # Own files + user/group/role shares. Public link rows (FileShare) must not widen discovery.
        scope = FileNode.query.filter(*base_filters, or_(*share_scope))

    limit = min(request.args.get("limit", default=1500, type=int), 5000)
    scan_cap = min(max(limit * 12, 500), 25000)
    candidates = scope.order_by(FileNode.is_folder.desc(), FileNode.name.asc()).limit(scan_cap).all()

    items: list[FileNode] = []
    for n in candidates:
        ok, _reason = access.can_access_node(current_user, n, "read")
        if not ok:
            continue
        if not access.documents_listing_includes_node(current_user, n):
            continue
        items.append(n)
        if len(items) >= limit:
            break

    truncated = len(items) >= limit or (len(candidates) >= scan_cap and len(items) > 0)

    ids = [n.id for n in items]
    shared_out_ids: set[int] = set()
    if ids:
        shared_out_ids |= _shared_out_node_ids(ids)

    payload_items = _serialize_nodes_batch(items, shared_out_ids=shared_out_ids)

    qdisp = raw[:80] + ("…" if len(raw) > 80 else "")
    crumbs = [{"id": None, "name": "All files"}, {"id": None, "name": f'Search: "{qdisp}"'}]
    default_parent = _default_home_folder_for_user(current_user.id)
    _audit("files.search", None, None, True, {"q": raw, "count": len(payload_items)})
    return jsonify(
        {
            "parent": None,
            "items": payload_items,
            "breadcrumb": crumbs,
            "shared_with_me": [],
            "shared_by_me": [],
            "default_parent_id": (default_parent.id if default_parent else None),
            "search": {"q": raw, "truncated": truncated, "tokens": tokens},
        }
    )


def _serialize_node(n: FileNode, include_children_count: bool = False) -> dict[str, Any]:
    data: dict[str, Any] = {
        "id": n.id,
        "name": n.name,
        "is_folder": n.is_folder,
        "parent_id": n.parent_id,
        "owner_id": n.owner_id,
        "created_at": n.created_at.isoformat() if n.created_at else None,
        "updated_at": n.updated_at.isoformat() if n.updated_at else None,
        "path_key": n.path_key,
        "attributes": n.attributes or {},
    }
    if not n.is_folder:
        cur = (
            db.session.query(FileVersion)
            .filter_by(file_node_id=n.id, is_current=True)
            .order_by(FileVersion.version_number.desc())
            .first()
        )
        if cur:
            data["size_bytes"] = cur.size_bytes
            data["mime_type"] = cur.mime_type
            data["version_number"] = cur.version_number
            data["sha256"] = cur.sha256
    if include_children_count and n.is_folder:
        data["child_count"] = n.children.count()
    if not n.is_folder:
        data["lock"] = lock_payload(get_lock_for_node(n.id))
        attach_edit_fields([data])
    if files_workspace.is_users_container_folder(n) and _shows_users_root_display_name(current_user):
        data["name"] = files_workspace.USERS_ROOT_DISPLAY_NAME
    return data


@bp.route("/api/upload-conflict", methods=["GET"])
@login_required
def api_upload_conflict():
    """Return whether a file with the same name already exists in the folder (for upload UI)."""
    parent_id = request.args.get("parent_id", type=int)
    filename = (request.args.get("filename") or "").strip()
    if parent_id is None or not filename:
        return jsonify({"error": "parent_id and filename required"}), 400
    name = filename.rsplit("/", 1)[-1]
    if not name or "/" in name or "\\" in name:
        return jsonify({"error": "invalid filename"}), 400
    parent = _node_or_404(parent_id)
    if not parent.is_folder:
        return jsonify({"error": "not a folder"}), 400
    ok, reason = access.can_access_node(current_user, parent, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    existing = (
        FileNode.query.filter_by(parent_id=parent.id, name=name, is_folder=False)
        .filter(FileNode.deleted_at.is_(None))
        .first()
    )
    if not existing:
        return jsonify({"conflict": False})

    ok_read, _ = access.can_access_node(current_user, existing, "read")
    if not ok_read:
        return jsonify({"conflict": False})

    can_replace, _ = access.can_access_node(current_user, existing, "write")
    cur = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=existing.id, is_current=True)
        .order_by(FileVersion.version_number.desc())
        .first()
    )
    return jsonify(
        {
            "conflict": True,
            "can_replace": can_replace,
            "existing": {
                "name": existing.name,
                "size_bytes": cur.size_bytes if cur else None,
                "updated_at": existing.updated_at.isoformat() if existing.updated_at else None,
                "mime_type": cur.mime_type if cur else None,
            },
        }
    )


@bp.route("/api/mkdir", methods=["POST"])
@login_required
def api_mkdir():
    payload = request.get_json(force=True, silent=True) or {}
    name = (payload.get("name") or "").strip()
    raw_parent_id = payload.get("parent_id")
    parent_id = int(raw_parent_id) if raw_parent_id is not None else None
    reuse_existing = bool(payload.get("reuse_existing"))
    if not name or "/" in name or "\\" in name:
        return jsonify({"error": "invalid name"}), 400
    parent = _resolve_parent_for_create(parent_id, current_user)
    if parent_id is None and parent is None and not _is_files_tree_admin(current_user):
        return jsonify({"error": "open a folder first (no folder found for your account)"}), 400
    if parent is not None:
        if not parent.is_folder:
            return jsonify({"error": "parent not folder"}), 400
        if files_workspace.destination_is_blocked_users_root(current_user, parent):
            return jsonify({"error": "forbidden", "reason": "cannot create in Users root"}), 403
        ok, reason = access.can_access_node(current_user, parent, "write")
        if not ok:
            _audit("files.mkdir.denied", "folder", str(parent.id), False, {"reason": reason})
            return jsonify({"error": "forbidden", "reason": reason}), 403
        if not _user_may_create_folder(current_user, parent):
            _audit("files.mkdir.denied", "folder", str(parent.id), False, {"reason": "create folders"})
            return jsonify({"error": "forbidden", "reason": "cannot create folders"}), 403

    owner_id = current_user.id
    if parent:
        owner_id = parent.owner_id

    if reuse_existing:
        # Folder upload wants to merge into an existing folder path (no auto-suffixing).
        existing_folder = (
            FileNode.query.filter_by(parent_id=(parent.id if parent else None), name=name, is_folder=True)
            .filter(FileNode.deleted_at.is_(None))
            .first()
        )
        if existing_folder:
            ok2, reason2 = access.can_access_node(current_user, existing_folder, "write")
            if not ok2:
                return jsonify({"error": "forbidden", "reason": reason2}), 403
            return jsonify({"node": _serialize_node(existing_folder)})
        # If there's an existing file with the same name, fail explicitly.
        existing_file = (
            FileNode.query.filter_by(parent_id=(parent.id if parent else None), name=name, is_folder=False)
            .filter(FileNode.deleted_at.is_(None))
            .first()
        )
        if existing_file:
            return jsonify({"error": "name exists (file)"}), 409

    # Avoid duplicate names: auto-suffix like "Name (1)".
    if parent:
        name = _next_available_name(parent.id, name)
    else:
        name = _next_available_name(None, name)

    node = FileNode(
        name=name,
        is_folder=True,
        parent_id=parent.id if parent else None,
        owner_id=owner_id,
        attributes=(payload.get("attributes") if isinstance(payload.get("attributes"), dict) else {}),
    )
    db.session.add(node)
    db.session.flush()
    node.path_key = _path_key_for(node)
    db.session.commit()
    _audit("files.mkdir", "folder", str(node.id), True, {"path": node.path_key})
    return jsonify({"node": _serialize_node(node)})


@bp.route("/api/upload", methods=["POST"])
@login_required
def api_upload():
    parent_id = request.form.get("parent_id", type=int)
    if parent_id is None:
        return jsonify({"error": "parent_id required"}), 400
    parent = _node_or_404(parent_id)
    if not parent.is_folder:
        return jsonify({"error": "parent not folder"}), 400
    ok, reason = access.can_access_node(current_user, parent, "write")
    if not ok:
        _audit("files.upload.denied", "folder", str(parent.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if files_workspace.destination_is_blocked_users_root(current_user, parent):
        return jsonify({"error": "forbidden", "reason": "cannot upload to users root"}), 403

    files = request.files.getlist("file")
    if not files:
        return jsonify({"error": "no file"}), 400

    results = []
    for f in files:
        if not f or not f.filename:
            continue
        relpath, size, sha256, mime = store_stream_and_digest(f.stream, f.filename)
        name = f.filename.rsplit("/", 1)[-1]
        existing = (
            FileNode.query.filter_by(parent_id=parent.id, name=name, is_folder=False)
            .filter(FileNode.deleted_at.is_(None))
            .first()
        )
        if existing:
            ok2, _ = access.can_access_node(current_user, existing, "write")
            if not ok2:
                results.append({"name": name, "error": "forbidden overwrite"})
                continue
            lock_ok, lock_err = assert_can_edit_locked_node(
                existing, current_user, files_admin=_is_files_tree_admin(current_user)
            )
            if not lock_ok:
                results.append({"name": name, "error": lock_err or "locked"})
                continue
            db.session.query(FileVersion).filter_by(file_node_id=existing.id, is_current=True).update({"is_current": False})
            last_v = (
                db.session.query(FileVersion)
                .filter_by(file_node_id=existing.id)
                .order_by(FileVersion.version_number.desc())
                .first()
            )
            vn = (last_v.version_number + 1) if last_v else 1
            ver = FileVersion(
                file_node_id=existing.id,
                version_number=vn,
                storage_relpath=relpath,
                size_bytes=size,
                sha256=sha256,
                mime_type=mime,
                created_by_id=current_user.id,
                is_current=True,
            )
            db.session.add(ver)
            existing.updated_at = utcnow()
            db.session.commit()
            _audit("files.upload.version", "file", str(existing.id), True, {"version": vn, "sha256": sha256})
            results.append({"name": name, "node": _serialize_node(existing), "version": vn})
            continue

        node = FileNode(
            name=name,
            is_folder=False,
            parent_id=parent.id,
            owner_id=parent.owner_id,
            attributes={},
        )
        db.session.add(node)
        db.session.flush()
        node.path_key = _path_key_for(node)
        ver = FileVersion(
            file_node_id=node.id,
            version_number=1,
            storage_relpath=relpath,
            size_bytes=size,
            sha256=sha256,
            mime_type=mime,
            created_by_id=current_user.id,
            is_current=True,
        )
        db.session.add(ver)
        db.session.commit()
        _audit("files.upload", "file", str(node.id), True, {"path": node.path_key, "sha256": sha256})
        results.append({"name": name, "node": _serialize_node(node), "version": 1})

    return jsonify({"results": results})


def _blank_docx_bytes() -> bytes:
    """Generate a minimal valid blank .docx (OOXML zip)."""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>
""",
        )
        z.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "word/document.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p/>
    <w:sectPr/>
  </w:body>
</w:document>
""",
        )
    return buf.getvalue()


def _blank_xlsx_bytes() -> bytes:
    """Generate a minimal valid blank .xlsx (OOXML zip)."""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>
""",
        )
        z.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "xl/workbook.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="Sheet1" sheetId="1" r:id="rId1"/>
  </sheets>
</workbook>
""",
        )
        z.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "xl/worksheets/sheet1.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <sheetData/>
</worksheet>
""",
        )
    return buf.getvalue()


def _blank_pptx_bytes() -> bytes:
    """Generate a minimal valid blank .pptx (OOXML zip)."""
    # OnlyOffice is stricter than most viewers; a super-minimal PPTX can fail to open.
    # Use python-pptx to generate a standards-compliant blank deck.
    try:
        from pptx import Presentation  # type: ignore

        buf = BytesIO()
        prs = Presentation()
        # Ensure at least one slide exists.
        try:
            layout = prs.slide_layouts[6]  # blank
        except Exception:
            layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        prs.save(buf)
        return buf.getvalue()
    except Exception:
        # Fallback to the previous minimal generator (best-effort).
        buf = BytesIO()
        with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
            z.writestr(
                "[Content_Types].xml",
                """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
</Types>
""",
            )
            z.writestr(
                "_rels/.rels",
                """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>
""",
            )
            z.writestr(
                "ppt/presentation.xml",
                """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:sldIdLst>
    <p:sldId id="256" r:id="rId1"/>
  </p:sldIdLst>
</p:presentation>
""",
            )
            z.writestr(
                "ppt/_rels/presentation.xml.rels",
                """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>
""",
            )
            z.writestr(
                "ppt/slides/slide1.xml",
                """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree/></p:cSld>
</p:sld>
""",
            )
        return buf.getvalue()


def _blank_vsdx_bytes() -> bytes:
    """Generate a minimal valid blank .vsdx (OPC zip)."""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/visio/document.xml" ContentType="application/vnd.ms-visio.drawing.main+xml"/>
  <Override PartName="/visio/pages/pages.xml" ContentType="application/vnd.ms-visio.pages+xml"/>
  <Override PartName="/visio/pages/page1.xml" ContentType="application/vnd.ms-visio.page+xml"/>
</Types>
""",
        )
        z.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.microsoft.com/visio/2010/relationships/document" Target="visio/document.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "visio/_rels/document.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.microsoft.com/visio/2010/relationships/pages" Target="pages/pages.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "visio/document.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<VisioDocument xmlns="http://schemas.microsoft.com/office/visio/2012/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <DocumentProperties/>
  <Pages r:id="rId1"/>
</VisioDocument>
""",
        )
        z.writestr(
            "visio/pages/_rels/pages.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.microsoft.com/visio/2010/relationships/page" Target="page1.xml"/>
</Relationships>
""",
        )
        z.writestr(
            "visio/pages/pages.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Pages xmlns="http://schemas.microsoft.com/office/visio/2012/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <Page ID="1" Name="Page-1" NameU="Page-1" IsCustomName="0" IsCustomNameU="0" Background="0" BackPage="0" ViewScale="1" ViewCenterX="4.25" ViewCenterY="5.5" r:id="rId1"/>
</Pages>
""",
        )
        z.writestr(
            "visio/pages/page1.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<PageContents xmlns="http://schemas.microsoft.com/office/visio/2012/main">
  <Shapes/>
</PageContents>
""",
        )
    return buf.getvalue()


def _blank_drawio_bytes() -> bytes:
    """Generate a blank .drawio diagram (XML container)."""
    # diagrams.net accepts an empty diagram and will initialize a blank page.
    xml = """<mxfile host="app.diagrams.net">
  <diagram name="Page-1"></diagram>
</mxfile>
"""
    return xml.encode("utf-8")


@bp.route("/api/new-docx", methods=["POST"])
@login_required
def api_new_docx():
    payload = request.get_json(force=True, silent=True) or {}
    parent_id = payload.get("parent_id", None)
    name = (payload.get("name") or "New document.docx").strip()
    if not name:
        name = "New document.docx"
    if "/" in name or "\\" in name:
        return jsonify({"error": "invalid name"}), 400
    if not name.lower().endswith(".docx"):
        name = name + ".docx"

    parent = None
    if parent_id is None:
        parent = _default_home_folder_for_user(current_user.id)
        if not parent:
            return jsonify({"error": "open a folder first (no folder found for your account)"}), 400
    else:
        parent = _node_or_404(int(parent_id))
    if not parent.is_folder:
        return jsonify({"error": "parent not folder"}), 400
    ok, reason = access.can_access_node(current_user, parent, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    name = _next_available_name(parent.id, name)

    doc_bytes = _blank_docx_bytes()
    relpath, size, sha256, mime = store_stream_and_digest(BytesIO(doc_bytes), name)
    node = FileNode(
        name=name,
        is_folder=False,
        parent_id=parent.id,
        owner_id=parent.owner_id,
        attributes={},
    )
    db.session.add(node)
    db.session.flush()
    node.path_key = _path_key_for(node)
    ver = FileVersion(
        file_node_id=node.id,
        version_number=1,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type=mime or "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(ver)
    db.session.commit()
    _audit("files.office.create", "file", str(node.id), True, {"path": node.path_key})
    return jsonify({"node": _serialize_node(node)}), 201


def _api_new_office_impl(payload: dict, *, kind: str):
    parent_id = payload.get("parent_id", None)
    name = (payload.get("name") or "").strip() or f"New {kind}.{kind}"
    if "/" in name or "\\" in name:
        return jsonify({"error": "invalid name"}), 400
    if not name.lower().endswith(f".{kind}"):
        name = name + f".{kind}"

    parent = None
    if parent_id is None:
        parent = _default_home_folder_for_user(current_user.id)
        if not parent:
            return jsonify({"error": "open a folder first (no folder found for your account)"}), 400
    else:
        parent = _node_or_404(int(parent_id))
    if not parent.is_folder:
        return jsonify({"error": "parent not folder"}), 400
    ok, reason = access.can_access_node(current_user, parent, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    name = _next_available_name(parent.id, name)

    if kind == "xlsx":
        blob = _blank_xlsx_bytes()
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif kind == "pptx":
        blob = _blank_pptx_bytes()
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        blob = _blank_docx_bytes()
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    relpath, size, sha256, _ = store_stream_and_digest(BytesIO(blob), name)
    node = FileNode(
        name=name,
        is_folder=False,
        parent_id=parent.id,
        owner_id=parent.owner_id,
        attributes={},
    )
    db.session.add(node)
    db.session.flush()
    node.path_key = _path_key_for(node)
    ver = FileVersion(
        file_node_id=node.id,
        version_number=1,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type=mime,
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(ver)
    db.session.commit()
    _audit("files.office.create", "file", str(node.id), True, {"path": node.path_key, "kind": kind})
    return jsonify({"node": _serialize_node(node)}), 201


@bp.route("/api/new-xlsx", methods=["POST"])
@login_required
def api_new_xlsx():
    payload = request.get_json(force=True, silent=True) or {}
    payload["name"] = (payload.get("name") or "New spreadsheet.xlsx")
    return _api_new_office_impl(payload, kind="xlsx")


@bp.route("/api/new-pptx", methods=["POST"])
@login_required
def api_new_pptx():
    payload = request.get_json(force=True, silent=True) or {}
    payload["name"] = (payload.get("name") or "New presentation.pptx")
    return _api_new_office_impl(payload, kind="pptx")


@bp.route("/api/new-vsdx", methods=["POST"])
@login_required
def api_new_vsdx():
    # Back-compat: historically this created .vsdx, but the UI uses diagrams.net which saves .drawio.
    # Keep endpoint name but create a .drawio file.
    payload = request.get_json(force=True, silent=True) or {}
    parent_id = payload.get("parent_id", None)
    name = (payload.get("name") or "New diagram.drawio").strip() or "New diagram.drawio"
    if "/" in name or "\\" in name:
        return jsonify({"error": "invalid name"}), 400
    if not name.lower().endswith(".drawio"):
        name = name + ".drawio"

    parent = None
    if parent_id is None:
        parent = _default_home_folder_for_user(current_user.id)
        if not parent:
            return jsonify({"error": "open a folder first (no folder found for your account)"}), 400
    else:
        parent = _node_or_404(int(parent_id))
    if not parent.is_folder:
        return jsonify({"error": "parent not folder"}), 400
    ok, reason = access.can_access_node(current_user, parent, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    name = _next_available_name(parent.id, name)
    blob = _blank_drawio_bytes()
    relpath, size, sha256, _ = store_stream_and_digest(BytesIO(blob), name)
    node = FileNode(
        name=name,
        is_folder=False,
        parent_id=parent.id,
        owner_id=parent.owner_id,
        attributes={},
    )
    db.session.add(node)
    db.session.flush()
    node.path_key = _path_key_for(node)
    ver = FileVersion(
        file_node_id=node.id,
        version_number=1,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type="application/xml",
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(ver)
    db.session.commit()
    _audit("files.drawio.create", "file", str(node.id), True, {"path": node.path_key})
    return jsonify({"node": _serialize_node(node)}), 201


@bp.route("/api/download/<int:node_id>", methods=["GET"])
@login_required
def api_download(node_id: int):
    node = _node_or_404(node_id)
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        _audit("files.download.denied", "file", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    version_id = request.args.get("version_id", type=int)
    if version_id:
        ver = db.session.get(FileVersion, version_id)
        if not ver or ver.file_node_id != node.id:
            abort(404)
    else:
        ver = (
            db.session.query(FileVersion)
            .filter_by(file_node_id=node.id, is_current=True)
            .order_by(FileVersion.version_number.desc())
            .first()
        )
    if not ver:
        abort(404)
    path = absolute_path(ver.storage_relpath)
    _audit("files.download", "file", str(node.id), True, {"version": ver.version_number})
    return send_file(path, as_attachment=True, download_name=node.name)


@bp.route("/api/view/<int:node_id>", methods=["GET"])
@login_required
def api_view(node_id: int):
    node = _node_or_404(node_id)
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        _audit("files.view.denied", "file", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    version_id = request.args.get("version_id", type=int)
    if version_id:
        ver = db.session.get(FileVersion, version_id)
        if not ver or ver.file_node_id != node.id:
            abort(404)
    else:
        ver = (
            db.session.query(FileVersion)
            .filter_by(file_node_id=node.id, is_current=True)
            .order_by(FileVersion.version_number.desc())
            .first()
        )
    if not ver:
        abort(404)

    path = absolute_path(ver.storage_relpath)
    mt, _enc = mimetypes.guess_type(node.name)
    if not mt:
        mt = "application/octet-stream"
    _audit("files.view", "file", str(node.id), True, {"version": ver.version_number, "mimetype": mt})
    return send_file(path, mimetype=mt, as_attachment=False, download_name=node.name)


@bp.route("/api/text/<int:node_id>", methods=["GET"])
@login_required
def api_text_preview(node_id: int):
    """Return a safe text preview for common config/code-like files.

    - Limits payload size
    - Rejects likely-binary data (NUL byte detection)
    """
    node = _node_or_404(node_id)
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        _audit("files.text.denied", "file", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    ver = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=node.id, is_current=True)
        .order_by(FileVersion.version_number.desc())
        .first()
    )
    if not ver:
        abort(404)
    path = absolute_path(ver.storage_relpath)

    limit = 350_000  # ~350KB preview
    try:
        with open(path, "rb") as f:
            data = f.read(limit + 1)
    except OSError:
        return jsonify({"error": "could not read file"}), 500

    truncated = len(data) > limit
    if truncated:
        data = data[:limit]

    # Quick binary detection
    if b"\x00" in data:
        return jsonify({"error": "binary", "message": "This file appears to be binary and cannot be rendered as text."}), 400

    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode(errors="replace")

    _audit("files.text.view", "file", str(node.id), True, {"version": ver.version_number, "truncated": truncated})
    return jsonify(
        {
            "ok": True,
            "id": node.id,
            "name": node.name,
            "text": text,
            "truncated": truncated,
            "bytes": ver.size_bytes or len(data),
            "version": ver.version_number,
        }
    )


def _safe_text(s: str | None, *, limit: int = 200_000) -> str:
    if not s:
        return ""
    out = str(s)
    if len(out) > limit:
        out = out[:limit] + "\n\n[truncated]"
    return out


def _eml_extract(path) -> dict[str, Any]:
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    headers = {
        "Subject": _safe_text(msg.get("subject")),
        "From": _safe_text(msg.get("from")),
        "To": _safe_text(msg.get("to")),
        "Cc": _safe_text(msg.get("cc")),
        "Date": _safe_text(msg.get("date")),
    }

    text_body = ""
    html_body = ""
    attachments: list[dict[str, Any]] = []

    def _att_kind(filename: str, content_type: str) -> str:
        fn = (filename or "").lower()
        ct = (content_type or "").lower()
        if fn.endswith(".pdf") or ct == "application/pdf":
            return "pdf"
        if fn.endswith((".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".odt", ".ods", ".odp", ".rtf")):
            return "office"
        if fn.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg")):
            return "image"
        if fn.endswith((".txt", ".md", ".csv", ".log")):
            return "text"
        return "file"

    try:
        att_idx = 0
        for part in msg.walk():
            if part.is_multipart():
                continue
            ctype = (part.get_content_type() or "").lower()
            disp = (part.get_content_disposition() or "").lower()
            filename = part.get_filename()

            # Attachments first (including inline with filenames)
            if disp == "attachment" or (filename and disp in ("attachment", "inline", "")):
                size_bytes = None
                try:
                    payload = part.get_payload(decode=True)
                    if isinstance(payload, (bytes, bytearray)):
                        size_bytes = len(payload)
                except Exception:
                    size_bytes = None
                attachments.append(
                    {
                        "index": att_idx,
                        "filename": _safe_text(filename) if filename else "",
                        "content_type": _safe_text(ctype),
                        "size_bytes": size_bytes,
                        "kind": _att_kind(filename or "", ctype),
                    }
                )
                att_idx += 1
                continue

            # Body candidates
            if ctype == "text/plain" and not text_body:
                try:
                    text_body = _safe_text(part.get_content())
                except Exception:
                    text_body = ""
            elif ctype == "text/html" and not html_body:
                try:
                    html_body = _safe_text(part.get_content())
                except Exception:
                    html_body = ""
    except Exception:
        pass

    # If the message isn't multipart, get_content() may be the body.
    if not text_body and not html_body:
        try:
            ctype = (msg.get_content_type() or "").lower()
            if ctype == "text/html":
                html_body = _safe_text(msg.get_content())
            else:
                text_body = _safe_text(msg.get_content())
        except Exception:
            pass

    # Build safe srcdoc: sandboxed iframe + CSP to block remote loads.
    html_srcdoc = ""
    if html_body:
        html_srcdoc = """<!doctype html>
<html><head>
  <meta charset="utf-8">
  <meta http-equiv="Content-Security-Policy" content="default-src 'none'; img-src data: cid:; style-src 'unsafe-inline';">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <style>
    body { margin: 12px; font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; color: #222; }
    img { max-width: 100%%; height: auto; }
    pre { white-space: pre-wrap; word-break: break-word; }
  </style>
</head><body>""" + html_body + "</body></html>"

    return {
        "headers": headers,
        "text_body": text_body,
        "html_srcdoc": html_srcdoc,
        "attachments": attachments[:200],
    }


def _eml_node_and_version(node_id: int) -> tuple[FileNode, FileVersion, Any]:
    node = _node_or_404(node_id)
    if node.is_folder:
        abort(404)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        _audit("files.eml.denied", "file", str(node.id), False, {"reason": reason})
        abort(403)

    ver = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=node.id, is_current=True)
        .order_by(FileVersion.version_number.desc())
        .first()
    )
    if not ver:
        abort(404)
    path = absolute_path(ver.storage_relpath)
    return node, ver, path


def _eml_find_attachment_bytes(path, index: int) -> tuple[bytes, str, str]:
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    cur = -1
    for part in msg.walk():
        if part.is_multipart():
            continue
        disp = (part.get_content_disposition() or "").lower()
        filename = part.get_filename()
        if disp == "attachment" or (filename and disp in ("attachment", "inline", "")):
            cur += 1
            if cur != index:
                continue
            ctype = (part.get_content_type() or "application/octet-stream").lower()
            data = part.get_payload(decode=True) or b""
            if not isinstance(data, (bytes, bytearray)):
                data = bytes(str(data), "utf-8", "ignore")
            name = filename or f"attachment-{index}"
            return bytes(data), name, ctype
    raise KeyError("attachment not found")


@bp.route("/eml/<int:node_id>", methods=["GET"])
@login_required
def eml_viewer(node_id: int):
    node, ver, path = _eml_node_and_version(node_id)
    parsed = _eml_extract(path)
    shell = (request.args.get("shell") or "").strip().lower()
    # In intranet shell, we want the menu/tab chrome visible.
    # If an old embed-style link is used, strip `embed=1`.
    if shell == "intranet" and (request.args.get("embed") or "").strip() == "1":
        args = dict(request.args)
        args.pop("embed", None)
        return redirect(url_for("files.eml_viewer", node_id=node_id, **args))
    tpl = "eml_viewer_intranet.html" if shell == "intranet" else "eml_viewer.html"
    _audit("files.eml.view", "file", str(node.id), True, {"version": ver.version_number})
    return render_template(
        tpl,
        node=node,
        version=ver,
        headers=parsed["headers"],
        text_body=parsed["text_body"],
        html_srcdoc=parsed["html_srcdoc"],
        attachments=parsed["attachments"],
    )


@bp.route("/eml/<int:node_id>/attachment/<int:att_index>", methods=["GET"])
@login_required
def eml_attachment(node_id: int, att_index: int):
    node, ver, path = _eml_node_and_version(node_id)
    try:
        data, filename, ctype = _eml_find_attachment_bytes(path, att_index)
    except KeyError:
        abort(404)

    inline = (request.args.get("inline") or "").strip().lower() in ("1", "true", "yes")
    mt = ctype or mimetypes.guess_type(filename)[0] or "application/octet-stream"
    _audit(
        "files.eml.attachment.view" if inline else "files.eml.attachment.download",
        "file",
        str(node.id),
        True,
        {"version": ver.version_number, "attachment_index": att_index, "content_type": mt},
    )
    return send_file(
        BytesIO(data),
        mimetype=mt,
        as_attachment=not inline,
        download_name=filename,
        max_age=0,
    )


@bp.route("/api/eml/<int:node_id>/extract-attachment", methods=["POST"])
@login_required
def api_eml_extract_attachment(node_id: int):
    node, ver, path = _eml_node_and_version(node_id)
    payload = request.get_json(force=True, silent=True) or {}
    try:
        att_index = int(payload.get("index"))
    except (TypeError, ValueError):
        return jsonify({"error": "index required"}), 400

    if node.parent_id is None:
        return jsonify({"error": "cannot extract attachment from a root file"}), 400
    parent = _node_or_404(int(node.parent_id))
    if not parent.is_folder:
        return jsonify({"error": "invalid parent"}), 400
    ok, reason = access.can_access_node(current_user, parent, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if not _user_may_create_folder(current_user, parent):
        return jsonify({"error": "forbidden", "reason": "cannot create folders"}), 403

    try:
        data, filename, ctype = _eml_find_attachment_bytes(path, att_index)
    except KeyError:
        return jsonify({"error": "not found"}), 404

    base_eml = node.name[:-4] if node.name.lower().endswith(".eml") else node.name
    folder_name = f"{base_eml} attachments".strip() or "Email attachments"
    folder_name = _next_available_name(parent.id, folder_name)
    folder = FileNode(
        name=folder_name,
        is_folder=True,
        parent_id=parent.id,
        owner_id=parent.owner_id,
        attributes={"eml_attachments_for": node.id},
        created_at=utcnow(),
        updated_at=utcnow(),
    )
    db.session.add(folder)
    db.session.flush()
    folder.path_key = _path_key_for(folder)
    db.session.add(folder)
    db.session.flush()

    safe_name = (filename or f"attachment-{att_index}").strip() or f"attachment-{att_index}"
    safe_name = _next_available_name(folder.id, safe_name)
    relpath, size, sha256, mime_guess = store_stream_and_digest(BytesIO(data), safe_name)
    file_node = FileNode(
        name=safe_name,
        is_folder=False,
        parent_id=folder.id,
        owner_id=folder.owner_id,
        attributes={"eml_source_node_id": node.id, "eml_attachment_index": att_index},
        created_at=utcnow(),
        updated_at=utcnow(),
    )
    db.session.add(file_node)
    db.session.flush()
    file_node.path_key = _path_key_for(file_node)
    ver2 = FileVersion(
        file_node_id=file_node.id,
        version_number=1,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type=(mime_guess or ctype or "application/octet-stream"),
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(ver2)
    db.session.commit()

    _audit(
        "files.eml.attachment.extract",
        "file",
        str(node.id),
        True,
        {"version": ver.version_number, "attachment_index": att_index, "new_node_id": file_node.id},
    )
    return jsonify({"node": _serialize_node(file_node)}), 201


@bp.route("/api/download-zip", methods=["POST"])
@login_required
def api_download_zip():
    payload = request.get_json(force=True, silent=True) or {}
    node_ids = payload.get("node_ids") or []
    if not isinstance(node_ids, list) or not node_ids:
        return jsonify({"error": "node_ids list required"}), 400

    # Deduplicate while preserving order.
    seen: set[int] = set()
    roots: list[FileNode] = []
    for raw in node_ids:
        try:
            nid = int(raw)
        except (TypeError, ValueError):
            continue
        if nid in seen:
            continue
        seen.add(nid)
        n = db.session.get(FileNode, nid)
        if not n:
            continue
        ok, reason = access.can_access_node(current_user, n, "read")
        if not ok:
            _audit("files.download.zip.denied", "node", str(nid), False, {"reason": reason})
            return jsonify({"error": "forbidden", "reason": reason}), 403
        roots.append(n)

    if not roots:
        return jsonify({"error": "no readable nodes"}), 400

    def safe_arcname(name: str) -> str:
        # Prevent zip-slip and normalize separators.
        s = (name or "item").replace("\\", "/")
        s = s.strip().lstrip("/").replace("../", "").replace("..", "")
        return s or "item"

    def rel_path_from_root(root: FileNode, node: FileNode) -> str:
        """Relative path like 'sub/child.txt' from root to node, without leading slash."""
        # Prefer path_key, falling back to walking parents.
        try:
            rk = (root.path_key or "").rstrip("/")
            nk = (node.path_key or "").rstrip("/")
            if rk and nk and nk.startswith(rk + "/"):
                rel = nk[len(rk) + 1 :]
                return safe_arcname(rel)
        except Exception:
            pass
        parts: list[str] = []
        cur: FileNode | None = node
        while cur is not None and cur.id != root.id:
            parts.append(safe_arcname(cur.name))
            cur = cur.parent
        parts.reverse()
        return "/".join([p for p in parts if p])

    buf = BytesIO()
    file_count = 0
    folder_count = 0
    download_name = "download.zip"
    if len(roots) == 1:
        download_name = f"{safe_arcname(roots[0].name)}.zip"
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
        for root in roots:
            base = safe_arcname(root.name)
            if root.is_folder:
                folder_count += 1
                # Add folder entries for nicer extraction.
                z.writestr(base.rstrip("/") + "/", b"")
                for n in _collect_subtree(root):
                    if n.id == root.id:
                        continue
                    rel = rel_path_from_root(root, n)
                    arc = base.rstrip("/") + "/" + rel
                    if n.is_folder:
                        z.writestr(arc.rstrip("/") + "/", b"")
                        continue
                    ok, reason = access.can_access_node(current_user, n, "read")
                    if not ok:
                        continue
                    ver = (
                        db.session.query(FileVersion)
                        .filter_by(file_node_id=n.id, is_current=True)
                        .order_by(FileVersion.version_number.desc())
                        .first()
                    )
                    if not ver:
                        continue
                    p = absolute_path(ver.storage_relpath)
                    if p.exists():
                        z.write(p, arcname=arc)
                        file_count += 1
            else:
                ok, reason = access.can_access_node(current_user, root, "read")
                if not ok:
                    continue
                ver = (
                    db.session.query(FileVersion)
                    .filter_by(file_node_id=root.id, is_current=True)
                    .order_by(FileVersion.version_number.desc())
                    .first()
                )
                if not ver:
                    continue
                p = absolute_path(ver.storage_relpath)
                if p.exists():
                    z.write(p, arcname=base)
                    file_count += 1

    buf.seek(0)
    _audit(
        "files.download.zip",
        "node",
        ",".join(str(n.id) for n in roots)[:200],
        True,
        {"files": file_count, "folders": folder_count, "roots": len(roots)},
    )
    return send_file(buf, mimetype="application/zip", as_attachment=True, download_name=download_name)


@bp.route("/api/favorites", methods=["GET"])
@login_required
def api_favorites_list():
    # List favorites for current user (readable nodes only).
    favs = (
        db.session.query(NodeFavorite)
        .filter_by(user_id=current_user.id)
        .order_by(NodeFavorite.created_at.desc())
        .all()
    )
    items: list[dict[str, Any]] = []
    for f in favs:
        n = db.session.get(FileNode, f.file_node_id)
        if not n:
            continue
        ok, _ = access.can_access_node(current_user, n, "read")
        if not ok:
            continue
        if not access.documents_listing_includes_node(current_user, n):
            continue
        d = _serialize_node(n, include_children_count=True)
        d["path_key"] = n.path_key or n.display_path()
        items.append(d)
    return jsonify({"items": items})


@bp.route("/api/favorites/<int:node_id>", methods=["PUT"])
@login_required
def api_favorite_put(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if not access.documents_listing_includes_node(current_user, node):
        return jsonify({"error": "forbidden", "reason": "not available in Documents"}), 403
    existing = NodeFavorite.query.filter_by(user_id=current_user.id, file_node_id=node.id).first()
    if existing:
        return jsonify({"ok": True, "favorite": True})
    db.session.add(NodeFavorite(user_id=current_user.id, file_node_id=node.id))
    db.session.commit()
    _audit("files.favorite.add", "node", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "favorite": True})


@bp.route("/api/favorites/<int:node_id>", methods=["DELETE"])
@login_required
def api_favorite_delete(node_id: int):
    node = _node_or_404(node_id)
    NodeFavorite.query.filter_by(user_id=current_user.id, file_node_id=node.id).delete(synchronize_session=False)
    db.session.commit()
    _audit("files.favorite.remove", "node", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "favorite": False})


@bp.route("/api/personal", methods=["GET"])
@login_required
def api_personal_list():
    # Personal files are an owner-only tag; return owned nodes marked personal.
    nodes = FileNode.query.filter_by(owner_id=current_user.id).order_by(FileNode.updated_at.desc()).all()
    items: list[dict[str, Any]] = []
    for n in nodes:
        if not (n.attributes or {}).get("personal"):
            continue
        if not access.documents_listing_includes_node(current_user, n):
            continue
        d = _serialize_node(n, include_children_count=True)
        d["path_key"] = n.path_key or n.display_path()
        items.append(d)
    return jsonify({"items": items})


@bp.route("/api/personal/<int:node_id>", methods=["PUT"])
@login_required
def api_personal_put(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "write" if node.is_folder else "read")
    if node.owner_id != current_user.id:
        return jsonify({"error": "forbidden", "reason": "only owner can mark personal"}), 403
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    attrs = dict(node.attributes or {})
    attrs["personal"] = True
    node.attributes = attrs
    node.updated_at = utcnow()
    db.session.add(node)
    db.session.commit()
    _audit("files.personal.add", "node", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "personal": True})


@bp.route("/api/personal/<int:node_id>", methods=["DELETE"])
@login_required
def api_personal_delete(node_id: int):
    node = _node_or_404(node_id)
    if node.owner_id != current_user.id:
        return jsonify({"error": "forbidden", "reason": "only owner can unmark personal"}), 403
    attrs = dict(node.attributes or {})
    attrs.pop("personal", None)
    node.attributes = attrs
    node.updated_at = utcnow()
    db.session.add(node)
    db.session.commit()
    _audit("files.personal.remove", "node", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "personal": False})


@bp.route("/api/move", methods=["PATCH"])
@login_required
def api_move():
    payload = request.get_json(force=True, silent=True) or {}
    node_id = payload.get("node_id")
    new_parent_id = payload.get("new_parent_id")
    if node_id is None or new_parent_id is None:
        return jsonify({"error": "node_id and new_parent_id required"}), 400
    node = _node_or_404(int(node_id))
    dest = _node_or_404(int(new_parent_id))
    if not dest.is_folder:
        return jsonify({"error": "destination not folder"}), 400
    if node.id == dest.id:
        return jsonify({"error": "invalid move"}), 400

    ok, reason = access.can_access_node(current_user, node, "move")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if not node.is_folder:
        lock_ok, lock_err = assert_can_edit_locked_node(
            node, current_user, files_admin=_is_files_tree_admin(current_user)
        )
        if not lock_ok:
            return jsonify({"error": "locked", "message": lock_err}), 423
    ok2, reason2 = access.can_access_node(current_user, dest, "write")
    if not ok2:
        return jsonify({"error": "forbidden", "reason": reason2}), 403
    if files_workspace.destination_is_blocked_users_root(current_user, dest):
        return jsonify({"error": "forbidden", "reason": "cannot move to users root"}), 403

    p = dest
    while p:
        if p.id == node.id:
            return jsonify({"error": "cannot move into self"}), 400
        p = p.parent

    if dest.owner_id != node.owner_id and not rbac.user_has_permission(current_user, rbac.PERMISSION_FILES_ADMIN):
        return jsonify({"error": "cross-owner move denied"}), 403

    clash = FileNode.query.filter_by(parent_id=dest.id, name=node.name).first()
    if clash:
        return jsonify({"error": "name exists in destination"}), 409

    node.parent_id = dest.id
    db.session.add(node)
    db.session.flush()
    subtree = _collect_subtree(node)
    for n in subtree:
        n.path_key = _path_key_for(n)
        db.session.add(n)
    db.session.commit()
    _audit("files.move", "node", str(node.id), True, {"to_parent": str(dest.id), "path": node.path_key})
    return jsonify({"node": _serialize_node(node)})


def _copy_node_recursive(src: FileNode, dest_parent: FileNode) -> FileNode:
    """Copy a node (file or folder) into dest_parent. Shares/favorites are not copied."""
    if not dest_parent.is_folder:
        raise ValueError("destination not folder")

    # Create new node
    new_node = FileNode(
        name=src.name,
        parent_id=dest_parent.id,
        is_folder=src.is_folder,
        owner_id=src.owner_id,
        attributes=dict(src.attributes or {}),
        created_at=utcnow(),
        updated_at=utcnow(),
    )
    db.session.add(new_node)
    db.session.flush()
    new_node.path_key = _path_key_for(new_node)
    db.session.add(new_node)
    db.session.flush()

    if src.is_folder:
        children = FileNode.query.filter_by(parent_id=src.id).order_by(FileNode.is_folder.desc(), FileNode.name.asc()).all()
        for ch in children:
            _copy_node_recursive(ch, new_node)
        return new_node

    # Copy file current version blob into new version 1
    ver = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=src.id, is_current=True)
        .order_by(FileVersion.version_number.desc())
        .first()
    )
    if not ver:
        return new_node
    relpath, size, sha256, mime = copy_blob_to_new_version(ver.storage_relpath)
    new_v = FileVersion(
        file_node_id=new_node.id,
        version_number=1,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type=mime or ver.mime_type,
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(new_v)
    return new_node


@bp.route("/api/copy", methods=["POST"])
@login_required
def api_copy():
    payload = request.get_json(force=True, silent=True) or {}
    node_id = payload.get("node_id")
    new_parent_id = payload.get("new_parent_id")
    if node_id is None or new_parent_id is None:
        return jsonify({"error": "node_id and new_parent_id required"}), 400
    node = _node_or_404(int(node_id))
    dest = _node_or_404(int(new_parent_id))
    if not dest.is_folder:
        return jsonify({"error": "destination not folder"}), 400

    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    ok2, reason2 = access.can_access_node(current_user, dest, "write")
    if not ok2:
        return jsonify({"error": "forbidden", "reason": reason2}), 403
    if files_workspace.destination_is_blocked_users_root(current_user, dest):
        return jsonify({"error": "forbidden", "reason": "cannot copy to users root"}), 403

    if dest.owner_id != node.owner_id and not rbac.user_has_permission(current_user, rbac.PERMISSION_FILES_ADMIN):
        return jsonify({"error": "cross-owner copy denied"}), 403

    clash = FileNode.query.filter_by(parent_id=dest.id, name=node.name).first()
    if clash:
        return jsonify({"error": "name exists in destination"}), 409

    new_node = _copy_node_recursive(node, dest)
    db.session.commit()
    _audit("files.copy", "node", str(node.id), True, {"to_parent": str(dest.id), "new_node_id": str(new_node.id), "path": new_node.path_key})
    return jsonify({"node": _serialize_node(new_node)}), 201


@bp.route("/api/node/<int:node_id>", methods=["PATCH"])
@login_required
def api_node_patch(node_id: int):
    """Rename a file or folder in place (same parent)."""
    node = _node_or_404(node_id)
    payload = request.get_json(force=True, silent=True) or {}
    if "name" not in payload:
        return jsonify({"error": "name required"}), 400
    name = (payload.get("name") or "").strip()
    if not name or "/" in name or "\\" in name:
        return jsonify({"error": "invalid name"}), 400
    if name == node.name:
        return jsonify({"node": _serialize_node(node)})

    ok, reason = access.can_access_node(current_user, node, "write")
    if not ok:
        _audit("files.rename.denied", "node", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    lock_ok, lock_err = assert_can_edit_locked_node(
        node, current_user, files_admin=_is_files_tree_admin(current_user)
    )
    if not lock_ok:
        return jsonify({"error": "locked", "message": lock_err}), 423

    # Ignore items in Recycle Bin; they should not block renames in the live tree.
    q = FileNode.query.filter_by(parent_id=node.parent_id, name=name).filter(FileNode.deleted_at.is_(None))
    clash = q.filter(FileNode.id != node.id).first()
    if clash:
        return jsonify({"error": "name already exists"}), 409

    old_path = node.path_key
    node.name = name
    node.updated_at = utcnow()
    db.session.add(node)
    db.session.flush()
    subtree = _collect_subtree(node)
    for n in subtree:
        n.path_key = _path_key_for(n)
        db.session.add(n)
    db.session.commit()
    _audit("files.rename", "node", str(node.id), True, {"from": old_path, "to": node.path_key})
    return jsonify({"node": _serialize_node(node)})


@bp.route("/api/lock/<int:node_id>", methods=["POST"])
@login_required
def api_lock_node(node_id: int):
    node = _node_or_404(node_id)
    if node.deleted_at is not None:
        return jsonify({"error": "not found"}), 404
    ok, reason = access.can_access_node(current_user, node, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    lock_data, err = acquire_lock(node, current_user)
    if err:
        code = 409 if lock_data else 400
        return jsonify({"error": err, "lock": lock_data}), code
    _audit("files.lock", "file", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "lock": lock_data, "node": _serialize_node(node)})


@bp.route("/api/unlock/<int:node_id>", methods=["POST"])
@login_required
def api_unlock_node(node_id: int):
    node = _node_or_404(node_id)
    if node.deleted_at is not None:
        return jsonify({"error": "not found"}), 404
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    released, err = release_lock(node, current_user, files_admin=_is_portal_admin(current_user))
    if err:
        return jsonify({"error": err}), 403
    if released:
        _audit("files.unlock", "file", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "node": _serialize_node(node)})


@bp.route("/api/edit-session/<int:node_id>", methods=["POST"])
@login_required
def api_edit_session_touch(node_id: int):
    node = _node_or_404(node_id)
    if node.deleted_at is not None:
        return jsonify({"error": "not found"}), 404
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    can_edit, _ = access.can_access_node(current_user, node, "write")
    if not can_edit:
        return jsonify({"error": "forbidden", "reason": "write required"}), 403
    row = touch_edit_session(node, current_user)
    if not row:
        return jsonify({"error": "could not register session"}), 400
    return jsonify({"ok": True, "session": row})


@bp.route("/api/edit-session/<int:node_id>", methods=["DELETE"])
@login_required
def api_edit_session_release(node_id: int):
    node = _node_or_404(node_id)
    if node.deleted_at is not None:
        return jsonify({"error": "not found"}), 404
    release_edit_session(node, current_user)
    return jsonify({"ok": True})


@bp.route("/api/edit-sessions", methods=["GET"])
@login_required
def api_edit_sessions():
    node_ids = request.args.getlist("node_id", type=int)
    node_ids = [int(n) for n in node_ids if n is not None]
    if not node_ids:
        return jsonify({"sessions": {}})
    node_ids = list(dict.fromkeys(node_ids))[:500]
    allowed: list[int] = []
    for nid in node_ids:
        node = db.session.get(FileNode, nid)
        if not node or node.deleted_at is not None or node.is_folder:
            continue
        ok, _ = access.can_access_node(current_user, node, "read")
        if ok:
            allowed.append(int(nid))
    return jsonify({"sessions": edit_sessions_map(allowed)})


@bp.route("/api/node/<int:node_id>", methods=["DELETE"])
@login_required
def api_delete(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "delete")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    if not node.is_folder:
        lock_ok, lock_err = assert_can_edit_locked_node(
            node, current_user, files_admin=_is_files_tree_admin(current_user)
        )
        if not lock_ok:
            return jsonify({"error": "locked", "message": lock_err}), 423

    payload = request.get_json(force=True, silent=True) or {}
    j_err, justification = validate_deletion_justification(payload)
    if j_err:
        return jsonify({"error": j_err}), 400

    path_snap = node.path_key
    try:
        marked = _soft_delete_node_recursive(node, deleted_by_id=current_user.id)
        db.session.commit()
    except SQLAlchemyError as exc:
        db.session.rollback()
        _audit("files.trash", "node", str(node_id), False, {"error": str(exc)})
        return jsonify({"error": "delete failed", "reason": str(exc)}), 500

    _audit(
        "files.trash",
        "node",
        str(node_id),
        True,
        {"path": path_snap, "count": marked, "justification": justification},
    )
    return jsonify({"ok": True, "trashed": True, "count": marked})


@bp.route("/api/versions/<int:node_id>", methods=["GET"])
@login_required
def api_versions(node_id: int):
    node = _node_or_404(node_id)
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "versions")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    vers = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=node.id)
        .order_by(FileVersion.version_number.desc())
        .all()
    )
    _audit("files.versions.list", "file", str(node.id), True, {"count": len(vers)})
    return jsonify(
        {
            "file": _serialize_node(node),
            "versions": [
                {
                    "id": v.id,
                    "version_number": v.version_number,
                    "size_bytes": v.size_bytes,
                    "sha256": v.sha256,
                    "created_at": v.created_at.isoformat(),
                    "created_by_id": v.created_by_id,
                    "is_current": v.is_current,
                }
                for v in vers
            ],
        }
    )


@bp.route("/api/restore", methods=["POST"])
@login_required
def api_restore():
    payload = request.get_json(force=True, silent=True) or {}
    node_id = payload.get("file_node_id")
    version_id = payload.get("version_id")
    if node_id is None or version_id is None:
        return jsonify({"error": "file_node_id and version_id required"}), 400
    node = _node_or_404(int(node_id))
    if node.is_folder:
        return jsonify({"error": "folder"}), 400
    ok, reason = access.can_access_node(current_user, node, "write")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    ver = db.session.get(FileVersion, int(version_id))
    if not ver or ver.file_node_id != node.id:
        abort(404)

    relpath, size, sha256, mime = copy_blob_to_new_version(ver.storage_relpath)
    db.session.query(FileVersion).filter_by(file_node_id=node.id, is_current=True).update({"is_current": False})
    last_v = (
        db.session.query(FileVersion)
        .filter_by(file_node_id=node.id)
        .order_by(FileVersion.version_number.desc())
        .first()
    )
    vn = (last_v.version_number + 1) if last_v else 1
    new_v = FileVersion(
        file_node_id=node.id,
        version_number=vn,
        storage_relpath=relpath,
        size_bytes=size,
        sha256=sha256,
        mime_type=mime or ver.mime_type,
        created_by_id=current_user.id,
        is_current=True,
    )
    db.session.add(new_v)
    node.updated_at = utcnow()
    db.session.commit()
    _audit("files.version.restore", "file", str(node.id), True, {"from_version": ver.version_number, "new_version": vn})
    return jsonify({"ok": True, "new_version": vn})


@bp.route("/api/track/<int:node_id>", methods=["GET"])
@login_required
def api_track(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    data: dict[str, Any] = {
        "tracking": {
            "id": node.id,
            "logical_path": node.path_key or node.display_path(),
            "owner_id": node.owner_id,
            "created_at": node.created_at.isoformat() if node.created_at else None,
            "updated_at": node.updated_at.isoformat() if node.updated_at else None,
            "attributes": node.attributes or {},
        }
    }
    if not node.is_folder:
        vers = (
            db.session.query(FileVersion)
            .filter_by(file_node_id=node.id)
            .order_by(FileVersion.version_number.desc())
            .all()
        )
        data["tracking"]["versions"] = [
            {
                "id": v.id,
                "version_number": v.version_number,
                "sha256": v.sha256,
                "size_bytes": v.size_bytes,
                "created_at": v.created_at.isoformat(),
                "is_current": v.is_current,
            }
            for v in vers
        ]
    _audit("files.track", "node", str(node.id), True, {})
    return jsonify(data)


def _serialize_internal_share(sh: NodeUserShare) -> dict[str, Any]:
    sw = sh.shared_with
    gb = sh.granted_by
    return {
        "id": sh.id,
        "share_type": "user",
        "username": sw.username if sw else "",
        "permission": sh.permission,
        "granted_by_username": gb.username if gb else "",
        "created_at": sh.created_at.isoformat() if sh.created_at else None,
    }


def _serialize_group_share(sh: NodeGroupShare) -> dict[str, Any]:
    grp = sh.group
    gb = sh.granted_by
    return {
        "id": sh.id,
        "share_type": "group",
        "group_name": grp.name if grp else "",
        "permission": sh.permission,
        "granted_by_username": gb.username if gb else "",
        "created_at": sh.created_at.isoformat() if sh.created_at else None,
    }


def _serialize_role_share(sh: NodeRoleShare) -> dict[str, Any]:
    role = sh.role
    gb = sh.granted_by
    return {
        "id": sh.id,
        "share_type": "role",
        "role_name": role.name if role else "",
        "permission": sh.permission,
        "granted_by_username": gb.username if gb else "",
        "created_at": sh.created_at.isoformat() if sh.created_at else None,
    }


def _list_internal_shares_for_node(node_id: int) -> list[dict[str, Any]]:
    rows: list[tuple[Any, dict[str, Any]]] = []
    for sh in NodeUserShare.query.filter_by(file_node_id=node_id).order_by(NodeUserShare.created_at.asc()):
        rows.append((sh.created_at, _serialize_internal_share(sh)))
    for sh in NodeGroupShare.query.filter_by(file_node_id=node_id).order_by(NodeGroupShare.created_at.asc()):
        rows.append((sh.created_at, _serialize_group_share(sh)))
    for sh in NodeRoleShare.query.filter_by(file_node_id=node_id).order_by(NodeRoleShare.created_at.asc()):
        rows.append((sh.created_at, _serialize_role_share(sh)))
    rows.sort(key=lambda pair: pair[0] or utcnow())
    return [pair[1] for pair in rows]


def _shared_out_node_ids(node_ids: list[int]) -> set[int]:
    if not node_ids:
        return set()
    out: set[int] = set()
    out |= {
        int(i)
        for (i,) in db.session.execute(
            select(NodeUserShare.file_node_id).where(NodeUserShare.file_node_id.in_(node_ids)).distinct()
        ).all()
    }
    out |= {
        int(i)
        for (i,) in db.session.execute(
            select(NodeGroupShare.file_node_id).where(NodeGroupShare.file_node_id.in_(node_ids)).distinct()
        ).all()
    }
    out |= {
        int(i)
        for (i,) in db.session.execute(
            select(NodeRoleShare.file_node_id).where(NodeRoleShare.file_node_id.in_(node_ids)).distinct()
        ).all()
    }
    out |= {
        int(i)
        for (i,) in db.session.execute(
            select(FileShare.file_node_id).where(FileShare.file_node_id.in_(node_ids)).distinct()
        ).all()
    }
    return out


def _group_role_share_path_prefixes(user: User) -> list[str]:
    """path_key prefixes for folders shared to this user via group or role membership."""
    prefixes: list[str] = []
    group_ids = {g.id for g in (user.groups or []) if g}
    if group_ids:
        for sh in NodeGroupShare.query.filter(NodeGroupShare.group_id.in_(group_ids)).all():
            n = sh.file_node
            if n and n.path_key:
                prefixes.append(n.path_key.rstrip("/"))
    role_ids = {
        int(rid)
        for (rid,) in db.session.execute(select(user_roles.c.role_id).where(user_roles.c.user_id == user.id)).all()
    }
    if role_ids:
        for sh in NodeRoleShare.query.filter(NodeRoleShare.role_id.in_(role_ids)).all():
            n = sh.file_node
            if n and n.path_key:
                prefixes.append(n.path_key.rstrip("/"))
    return list(dict.fromkeys(p for p in prefixes if p))


def _append_shared_with_me_entries(user: User, out: list[dict[str, Any]], seen: set[int]) -> None:
    """Add nodes the user can access via explicit or group/role shares (not owned)."""
    candidates: list[tuple[int, str | None]] = []
    for sh in NodeUserShare.query.filter_by(shared_with_user_id=user.id).all():
        candidates.append((sh.file_node_id, sh.permission))
    group_ids = {g.id for g in (user.groups or []) if g}
    if group_ids:
        for sh in NodeGroupShare.query.filter(NodeGroupShare.group_id.in_(group_ids)).all():
            candidates.append((sh.file_node_id, sh.permission))
    role_ids = {
        int(rid)
        for (rid,) in db.session.execute(select(user_roles.c.role_id).where(user_roles.c.user_id == user.id)).all()
    }
    if role_ids:
        for sh in NodeRoleShare.query.filter(NodeRoleShare.role_id.in_(role_ids)).all():
            candidates.append((sh.file_node_id, sh.permission))

    pending: list[tuple[FileNode, str]] = []
    node_ids_to_fetch: list[int] = []
    for node_id, _hint_perm in candidates:
        if node_id in seen:
            continue
        seen.add(node_id)
        node_ids_to_fetch.append(node_id)

    if not node_ids_to_fetch:
        return
    nodes_by_id = {
        n.id: n
        for n in FileNode.query.filter(FileNode.id.in_(node_ids_to_fetch), FileNode.deleted_at.is_(None)).all()
    }
    for node_id in node_ids_to_fetch:
        sn = nodes_by_id.get(node_id)
        if not sn or sn.owner_id == user.id:
            continue
        ok, _ = access.can_access_node(user, sn, "read")
        if not ok or not access.documents_listing_includes_node(user, sn):
            continue
        perm = access.internal_share_grant(user.id, sn) or "read"
        pending.append((sn, perm))

    if not pending:
        return
    batch = _serialize_nodes_batch([sn for sn, _ in pending], include_owner_username=True)
    for d, (sn, perm) in zip(batch, pending):
        d["path_key"] = sn.path_key or sn.display_path()
        d["permission"] = perm
        out.append(d)


@bp.route("/api/node/<int:node_id>/detail", methods=["GET"])
@login_required
def api_node_detail(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    role = access.effective_access_role(current_user, node)
    owner = db.session.get(User, node.owner_id)
    owner_username = owner.username if owner else ""
    link_q = db.session.query(FileShare).filter_by(file_node_id=node.id)
    link_count = link_q.count()
    can_manage = access.can_manage_user_shares(current_user, node)
    link_shares: list[dict[str, Any]] = []
    if can_manage and link_count:
        # Return active public-link URLs so UI can display/copy them.
        for sh in link_q.order_by(FileShare.id.desc()).limit(20).all():
            try:
                url = url_for("shares.public_share", token=sh.token, _external=True)
            except Exception:
                url = ""
            link_shares.append(
                {
                    "id": sh.id,
                    "token": sh.token,
                    "url": url,
                    "permission": sh.permission,
                    "expires_at": (sh.expires_at.isoformat() if sh.expires_at else None),
                }
            )

    internal_shares: list[dict[str, Any]] = []
    your_share: dict[str, Any] | None = None
    if can_manage:
        internal_shares = _list_internal_shares_for_node(node.id)
    elif role in ("read", "write"):
        perm = access.internal_share_grant(current_user.id, node)
        if perm:
            your_share = {"permission": perm}

    ndata = _serialize_node(node, include_children_count=bool(node.is_folder))
    ndata["owner_username"] = owner_username

    return jsonify(
        {
            "node": ndata,
            "role": role,
            "owner_username": owner_username,
            "can_manage_sharing": can_manage,
            "internal_shares": internal_shares,
            "your_share": your_share,
            "link_shares_count": link_count,
            "link_shares": link_shares,
        }
    )


@bp.route("/api/node/<int:node_id>/activity", methods=["GET"])
@login_required
def api_node_activity(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    limit = min(request.args.get("limit", default=50, type=int), 200)
    allowed_actions = [
        # "opened"
        "files.open",
        # changes
        "files.mkdir",
        "files.upload",
        "files.upload.version",
        "files.rename",
        "files.move",
        "files.delete",
        "files.version.restore",
        # sharing + chat
        "files.share.user.grant",
        "files.share.user.revoke",
        "files.comment.post",
        # access
        "files.download",
    ]
    q = (
        db.session.query(AuditLog)
        .filter(
            AuditLog.resource_id == str(node.id),
            AuditLog.resource_type.in_(["node", "file", "folder"]),
            AuditLog.action.in_(allowed_actions),
        )
        .order_by(AuditLog.timestamp.desc())
        .limit(limit)
    )
    rows = q.all()
    return jsonify(
        {
            "items": [
                {
                    "id": r.id,
                    "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                    "username": r.username_snapshot,
                    "action": r.action,
                    "success": r.success,
                    "details": r.details,
                }
                for r in rows
            ]
        }
    )


@bp.route("/api/node/<int:node_id>/comments", methods=["GET"])
@login_required
def api_node_comments(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    limit = min(request.args.get("limit", default=50, type=int), 200)
    rows = (
        FileComment.query.filter_by(file_node_id=node.id)
        .order_by(FileComment.created_at.desc())
        .limit(limit)
        .all()
    )
    return jsonify(
        {
            "items": [
                {
                    "id": c.id,
                    "created_at": c.created_at.isoformat() if c.created_at else None,
                    "user_id": c.user_id,
                    "username": c.user.username if c.user else None,
                    "body": c.body,
                }
                for c in rows
            ]
        }
    )


@bp.route("/api/node/<int:node_id>/comments", methods=["POST"])
@login_required
def api_post_node_comment(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "read")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    payload = request.get_json(force=True, silent=True) or {}
    body = (payload.get("body") or "").strip()
    if not body:
        return jsonify({"error": "comment required"}), 400
    if len(body) > 4000:
        return jsonify({"error": "comment too long"}), 400
    c = FileComment(file_node_id=node.id, user_id=current_user.id, body=body)
    db.session.add(c)
    db.session.commit()
    _audit("files.comment.post", "node", str(node.id), True, {"comment_id": c.id})
    return jsonify({"ok": True, "id": c.id}), 201


@bp.route("/api/users/suggest", methods=["GET"])
@login_required
def api_users_suggest():
    q = (request.args.get("q") or "").strip()
    if len(q) < 1:
        return jsonify({"users": []})
    like = f"%{q}%"
    users = (
        User.query.filter(
            User.is_active.is_(True),
            User.id != current_user.id,
            or_(User.username.ilike(like), func.coalesce(User.full_name, "").ilike(like), func.coalesce(User.email, "").ilike(like)),
        )
        .order_by(User.username)
        .limit(20)
        .all()
    )
    return jsonify(
        {
            "users": [
                {"id": u.id, "username": u.username, "full_name": (u.full_name or ""), "email": (u.email or "")}
                for u in users
            ]
        }
    )


@bp.route("/api/groups/suggest", methods=["GET"])
@login_required
def api_groups_suggest():
    q = (request.args.get("q") or "").strip()
    if len(q) < 1:
        return jsonify({"groups": []})
    like = f"%{q}%"
    groups = (
        Group.query.filter(or_(Group.name.ilike(like), func.coalesce(Group.description, "").ilike(like)))
        .order_by(Group.name)
        .limit(20)
        .all()
    )
    return jsonify(
        {
            "groups": [
                {
                    "id": g.id,
                    "name": g.name,
                    "description": (g.description or ""),
                    "member_count": len(g.users or []),
                }
                for g in groups
            ]
        }
    )


@bp.route("/api/roles/suggest", methods=["GET"])
@login_required
def api_roles_suggest():
    q = (request.args.get("q") or "").strip()
    if len(q) < 1:
        return jsonify({"roles": []})
    like = f"%{q}%"
    roles = Role.query.filter(Role.name.ilike(like)).order_by(Role.name).limit(20).all()
    return jsonify({"roles": [{"id": r.id, "name": r.name} for r in roles]})


@bp.route("/api/node/<int:node_id>/shares/user", methods=["POST"])
@login_required
def api_add_node_user_share(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        _audit("files.share.user.denied", "node", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    payload = request.get_json(force=True, silent=True) or {}
    username = (payload.get("username") or "").strip()
    permission = (payload.get("permission") or "read").lower()
    if permission not in ("read", "write"):
        return jsonify({"error": "invalid permission"}), 400
    if not username:
        return jsonify({"error": "username required"}), 400

    # UI accepts either username or email address.
    ident = username.lower()
    target = (
        User.query.filter(
            or_(
                db.func.lower(User.username) == ident,
                db.func.lower(func.coalesce(User.email, "")) == ident,
            )
        )
        .first()
    )
    if not target:
        return jsonify({"error": "user not found"}), 404
    if target.id == current_user.id:
        return jsonify({"error": "cannot share with yourself"}), 400
    if target.id == node.owner_id:
        return jsonify({"error": "owner already has access"}), 400

    existing = NodeUserShare.query.filter_by(file_node_id=node.id, shared_with_user_id=target.id).first()
    if existing:
        existing.permission = permission
        existing.granted_by_id = current_user.id
    else:
        existing = NodeUserShare(
            file_node_id=node.id,
            shared_with_user_id=target.id,
            permission=permission,
            granted_by_id=current_user.id,
        )
        db.session.add(existing)
    db.session.commit()
    db.session.refresh(existing)
    _audit(
        "files.share.user.upsert",
        "node",
        str(node.id),
        True,
        {"target": target.username, "permission": permission},
    )
    return jsonify({"share": _serialize_internal_share(existing)})


@bp.route("/api/node/<int:node_id>/shares/group", methods=["POST"])
@login_required
def api_add_node_group_share(node_id: int):
    """Share with a group tag; members inherit access dynamically."""
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        _audit("files.share.group.denied", "node", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    payload = request.get_json(force=True, silent=True) or {}
    group_name = (payload.get("group_name") or "").strip()
    permission = (payload.get("permission") or "read").lower()
    if permission not in ("read", "write"):
        return jsonify({"error": "invalid permission"}), 400
    if not group_name:
        return jsonify({"error": "group_name required"}), 400

    grp = Group.query.filter(func.lower(Group.name) == group_name.lower()).first()
    if not grp:
        return jsonify({"error": "group not found"}), 404

    existing = NodeGroupShare.query.filter_by(file_node_id=node.id, group_id=grp.id).first()
    if existing:
        existing.permission = permission
        existing.granted_by_id = current_user.id
    else:
        existing = NodeGroupShare(
            file_node_id=node.id,
            group_id=grp.id,
            permission=permission,
            granted_by_id=current_user.id,
        )
        db.session.add(existing)
    db.session.commit()
    db.session.refresh(existing)
    _audit(
        "files.share.group.upsert",
        "node",
        str(node.id),
        True,
        {"group": grp.name, "permission": permission},
    )
    return jsonify({"share": _serialize_group_share(existing)})


@bp.route("/api/node/<int:node_id>/shares/role", methods=["POST"])
@login_required
def api_add_node_role_share(node_id: int):
    """Share with a role tag; users with that role inherit access dynamically."""
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        _audit("files.share.role.denied", "node", str(node.id), False, {"reason": reason})
        return jsonify({"error": "forbidden", "reason": reason}), 403

    payload = request.get_json(force=True, silent=True) or {}
    role_name = (payload.get("role_name") or "").strip()
    permission = (payload.get("permission") or "read").lower()
    if permission not in ("read", "write"):
        return jsonify({"error": "invalid permission"}), 400
    if not role_name:
        return jsonify({"error": "role_name required"}), 400

    role = Role.query.filter(func.lower(Role.name) == role_name.lower()).first()
    if not role:
        return jsonify({"error": "role not found"}), 404

    existing = NodeRoleShare.query.filter_by(file_node_id=node.id, role_id=role.id).first()
    if existing:
        existing.permission = permission
        existing.granted_by_id = current_user.id
    else:
        existing = NodeRoleShare(
            file_node_id=node.id,
            role_id=role.id,
            permission=permission,
            granted_by_id=current_user.id,
        )
        db.session.add(existing)
    db.session.commit()
    db.session.refresh(existing)
    _audit(
        "files.share.role.upsert",
        "node",
        str(node.id),
        True,
        {"role": role.name, "permission": permission},
    )
    return jsonify({"share": _serialize_role_share(existing)})


@bp.route("/api/node/<int:node_id>/shares/user/<int:share_id>", methods=["DELETE"])
@login_required
def api_remove_node_user_share(node_id: int, share_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    sh = db.session.get(NodeUserShare, share_id)
    if not sh or sh.file_node_id != node.id:
        abort(404)
    db.session.delete(sh)
    db.session.commit()
    _audit("files.share.user.revoke", "node", str(node.id), True, {"share_id": share_id})
    return jsonify({"ok": True})


@bp.route("/api/node/<int:node_id>/shares/group/<int:share_id>", methods=["DELETE"])
@login_required
def api_remove_node_group_share(node_id: int, share_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    sh = db.session.get(NodeGroupShare, share_id)
    if not sh or sh.file_node_id != node.id:
        abort(404)
    db.session.delete(sh)
    db.session.commit()
    _audit("files.share.group.revoke", "node", str(node.id), True, {"share_id": share_id})
    return jsonify({"ok": True})


@bp.route("/api/node/<int:node_id>/shares/role/<int:share_id>", methods=["DELETE"])
@login_required
def api_remove_node_role_share(node_id: int, share_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403

    sh = db.session.get(NodeRoleShare, share_id)
    if not sh or sh.file_node_id != node.id:
        abort(404)
    db.session.delete(sh)
    db.session.commit()
    _audit("files.share.role.revoke", "node", str(node.id), True, {"share_id": share_id})
    return jsonify({"ok": True})


@bp.route("/api/node/<int:node_id>/shares/link", methods=["DELETE"])
@login_required
def api_remove_node_link_shares(node_id: int):
    node = _node_or_404(node_id)
    ok, reason = access.can_access_node(current_user, node, "share")
    if not ok:
        return jsonify({"error": "forbidden", "reason": reason}), 403
    n = FileShare.query.filter_by(file_node_id=node.id).delete(synchronize_session=False)
    db.session.commit()
    _audit("files.share.link.revoke", "node", str(node.id), True, {"count": n})
    return jsonify({"ok": True, "count": n})


def _recycle_retention_days() -> int:
    cfg = get_setting("recycle", default={}) or {}
    try:
        d = int(cfg.get("retention_days") or 1)
    except (TypeError, ValueError):
        d = 1
    return max(0, min(d, 3650))


def _purge_expired_deleted_for_owner(owner_id: int) -> int:
    days = _recycle_retention_days()
    if days <= 0:
        return 0
    cutoff = utcnow() - timedelta(days=days)
    rows = (
        FileNode.query.filter_by(owner_id=owner_id)
        .filter(FileNode.deleted_at.isnot(None))
        .filter(FileNode.deleted_at < cutoff)
        .order_by(FileNode.deleted_at.asc())
        .all()
    )
    purged = 0
    for n in rows:
        try:
            _delete_node_recursive(n)
            db.session.flush()
            purged += 1
        except SQLAlchemyError:
            db.session.rollback()
            break
    if purged:
        db.session.commit()
    return purged


@bp.route("/api/recycle", methods=["GET"])
@login_required
def api_recycle_list():
    _purge_expired_deleted_for_owner(current_user.id)
    limit = min(request.args.get("limit", default=200, type=int), 1000)
    rows = (
        FileNode.query.filter_by(owner_id=current_user.id)
        .filter(FileNode.deleted_at.isnot(None))
        .order_by(FileNode.deleted_at.desc())
        .limit(limit)
        .all()
    )
    items: list[dict[str, Any]] = []
    for n in rows:
        # Only show top-level deleted items (hide children of deleted folders).
        if n.parent_id is not None:
            p = db.session.get(FileNode, int(n.parent_id))
            if p and p.deleted_at is not None:
                continue
        d = _serialize_node(n, include_children_count=True)
        d["deleted_at"] = n.deleted_at.isoformat() if n.deleted_at else None
        d["deleted_by_id"] = n.deleted_by_id
        d["original_parent_id"] = n.original_parent_id
        items.append(d)
    return jsonify({"items": items, "retention_days": _recycle_retention_days()})


@bp.route("/api/recycle/<int:node_id>/restore", methods=["POST"])
@login_required
def api_recycle_restore(node_id: int):
    node = _node_or_404(node_id)
    if node.owner_id != current_user.id:
        return jsonify({"error": "forbidden", "reason": "only owner can restore"}), 403
    if node.deleted_at is None:
        return jsonify({"error": "not deleted"}), 400

    dest_parent = None
    if node.original_parent_id is not None:
        dest_parent = db.session.get(FileNode, int(node.original_parent_id))
    if not dest_parent or not dest_parent.is_folder or dest_parent.deleted_at is not None:
        dest_parent = _default_home_folder_for_user(current_user.id)
    if not dest_parent:
        return jsonify({"error": "no destination folder"}), 400
    if files_workspace.destination_is_blocked_users_root(current_user, dest_parent):
        dest_parent = _default_home_folder_for_user(current_user.id)
    if not dest_parent:
        return jsonify({"error": "no destination folder"}), 400

    # Avoid clashes by suffixing.
    new_name = _next_available_name(dest_parent.id, node.name)
    node.name = new_name
    node.parent_id = dest_parent.id
    node.deleted_at = None
    node.deleted_by_id = None
    node.original_parent_id = None
    db.session.add(node)
    db.session.flush()

    # Restore descendants and their parent links remain valid.
    for n in _collect_subtree(node):
        n.deleted_at = None
        n.deleted_by_id = None
        n.original_parent_id = None
        n.path_key = _path_key_for(n)
        db.session.add(n)
    db.session.commit()
    _audit("files.recycle.restore", "node", str(node.id), True, {"path": node.path_key})
    return jsonify({"ok": True, "node": _serialize_node(node)})


@bp.route("/api/recycle/<int:node_id>/purge", methods=["DELETE"])
@login_required
def api_recycle_purge(node_id: int):
    node = _node_or_404(node_id)
    if node.owner_id != current_user.id and not rbac.user_has_permission(current_user, rbac.PERMISSION_FILES_ADMIN):
        return jsonify({"error": "forbidden"}), 403
    if node.deleted_at is None:
        return jsonify({"error": "not deleted"}), 400

    payload = request.get_json(force=True, silent=True) or {}
    j_err, justification = validate_deletion_justification(payload)
    if j_err:
        return jsonify({"error": j_err}), 400

    path_snap = node.path_key
    try:
        _delete_node_recursive(node)
        db.session.commit()
    except SQLAlchemyError as exc:
        db.session.rollback()
        return jsonify({"error": "purge failed", "reason": str(exc)}), 500
    _audit(
        "files.recycle.purge",
        "node",
        str(node_id),
        True,
        {"path": path_snap, "justification": justification},
    )
    return jsonify({"ok": True})
